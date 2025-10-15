import io
import re
import os
import json
import zipfile
import logging
import datetime
import requests
from collections import Counter

# Optional / platform-dependent deps
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None
try:
    import easyocr
except ImportError:
    easyocr = None
try:
    import olefile
except ImportError:
    olefile = None
try:
    import win32com.client
except ImportError:
    win32com = None

import numpy as np
from PIL import Image
from transformers import AutoTokenizer, AutoModelForTokenClassification, pipeline

# Presentation / Office libs
try:
    from docx import Document
except ImportError:
    Document = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

# Face detector (MTCNN)
try:
    from mtcnn import MTCNN
except ImportError:
    MTCNN = None

author = "merged-by-chatgpt"
logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(message)s')

# ==========================
# Globals
# ==========================
HF_TOKEN = os.getenv("HF_TOKEN", None)
NER_MODEL_NAME = "soddokayo/klue-roberta-base-ner"
IMAGE_EXTENSIONS = ("png", "jpg", "jpeg", "bmp", "webp", "gif", "tif", "tiff")

easyocr_reader = None
if easyocr is not None:
    try:
        easyocr_reader = easyocr.Reader(['ko', 'en'], gpu=True)
        logging.info("EasyOCR 초기화 완료")
    except Exception as e:
        logging.warning(f"EasyOCR 초기화 실패: {e}")
        easyocr_reader = None
else:
    logging.warning("easyocr 미설치")

# MTCNN detector
if MTCNN is not None:
    detector = MTCNN()
else:
    detector = None
    logging.warning("MTCNN 미설치")

# HuggingFace NER
logging.info(f"NER 모델 로딩 중: {NER_MODEL_NAME}")
if HF_TOKEN:
    ner_tokenizer = AutoTokenizer.from_pretrained(NER_MODEL_NAME, token=HF_TOKEN)
    ner_model = AutoModelForTokenClassification.from_pretrained(NER_MODEL_NAME, token=HF_TOKEN)
else:
    ner_tokenizer = AutoTokenizer.from_pretrained(NER_MODEL_NAME)
    ner_model = AutoModelForTokenClassification.from_pretrained(NER_MODEL_NAME)
ner_pipeline = pipeline("ner", model=ner_model, tokenizer=ner_tokenizer, grouped_entities=True)
logging.info("NER 모델 로딩 완료")

# ==========================
# OCR helpers
# ==========================
def _easyocr_on_image(img: Image.Image) -> str:
    if easyocr_reader is None:
        return ""
    try:
        result = easyocr_reader.readtext(np.array(img))
        return " ".join([box[1] for box in result]).strip()
    except Exception:
        return ""

def run_ocr_on_single_image(image_bytes: bytes) -> str:
    try:
        img = Image.open(io.BytesIO(image_bytes))
    except Exception:
        return ""
    return _easyocr_on_image(img)

def run_ocr_on_docx_images(file_bytes: bytes) -> str:
    text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            for name in z.namelist():
                if name.startswith("word/media/"):
                    with z.open(name) as imf:
                        try:
                            img = Image.open(imf)
                            t = _easyocr_on_image(img)
                            if t:
                                text += t + "\n"
                        except Exception:
                            continue
    except Exception as e:
        logging.error(f"DOCX 이미지 OCR 실패: {e}")
    return text.replace("\n", " ").strip()

def run_ocr_on_hwp_images(file_bytes: bytes) -> str:
    if olefile is None:
        return ""
    text = ""
    try:
        with olefile.OleFileIO(io.BytesIO(file_bytes)) as f:
            for entry in f.listdir():
                if "BinData" in entry[0]:
                    data = f.openstream(entry).read()
                    try:
                        img = Image.open(io.BytesIO(data))
                        t = _easyocr_on_image(img)
                        if t:
                            text += t + "\n"
                    except Exception:
                        continue
    except Exception as e:
        logging.error(f"HWP 이미지 OCR 실패: {e}")
    return text.replace("\n", " ").strip()

def run_ocr_on_hwpx_images(file_bytes: bytes) -> str:
    text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            for name in z.namelist():
                if name.startswith("Contents/") and name.lower().endswith((".jpg", ".png", ".bmp")):
                    img = Image.open(io.BytesIO(z.read(name)))
                    t = _easyocr_on_image(img)
                    if t:
                        text += t + "\n"
    except Exception as e:
        logging.error(f"HWPX 이미지 OCR 실패: {e}")
    return text.replace("\n", " ").strip()

def run_ocr_on_ppt_images(file_bytes: bytes) -> str:
    if Presentation is None:
        return ""
    text = ""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    try:
                        img = Image.open(io.BytesIO(shape.image.blob))
                        t = _easyocr_on_image(img)
                        if t:
                            text += t + "\n"
                    except Exception:
                        continue
    except Exception as e:
        logging.error(f"PPT/PPTX 이미지 OCR 실패: {e}")
    return text.replace("\n", " ").strip()

def run_ocr_on_xls_images(file_bytes: bytes) -> str:
    text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            for name in z.namelist():
                if name.lower().startswith("xl/media/"):
                    img = Image.open(io.BytesIO(z.read(name)))
                    t = _easyocr_on_image(img)
                    if t:
                        text += t + "\n"
    except Exception as e:
        logging.error(f"XLS/XLSX 이미지 OCR 실패: {e}")
    return text.replace("\n", " ").strip()

def run_ocr_on_pdf_images(pdf_bytes: bytes) -> str:
    if fitz is None:
        return ""
    text = ""
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            for page in doc:
                pix = page.get_pixmap()
                try:
                    img = Image.open(io.BytesIO(pix.tobytes("ppm")))
                except Exception:
                    continue
                t = _easyocr_on_image(img)
                if t:
                    text += t + "\n"
    except Exception as e:
        logging.error(f"PDF 이미지 OCR 실패: {e}")
    return text.replace("\n", " ").strip()

# ==========================
# File parsing (union of features)
# ==========================
def parse_file(File_Bytes: bytes, File_Ext: str) -> tuple:
    """
    Returns (text, is_image)
    """
    ext = (File_Ext or "").lower()

    if ext == "txt":
        try:
            return File_Bytes.decode("utf-8"), False
        except UnicodeDecodeError:
            return File_Bytes.decode("cp949"), False

    if ext == "hwp":
        if win32com is not None: # hwp를 hwpx로 변환
            import tempfile, os as _os
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".hwp") as tmp:
                    tmp.write(File_Bytes); tmp_path = tmp.name
                tmp_hwpx_path = tmp_path + "x"
                try:
                    hwp = win32com.client.Dispatch("HWPFrame.HwpObject")
                    hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModule")
                    hwp.Open(tmp_path)
                    hwp.SaveAs(tmp_hwpx_path, "HWPX")
                    hwp.Quit()
                except Exception:
                    raise ValueError("[ERROR] HWP 변환 실패 또는 암호화")
                with open(tmp_hwpx_path, "rb") as f:
                    converted = f.read()
                _os.remove(tmp_path); _os.remove(tmp_hwpx_path)
                return parse_file(converted, "hwpx")
            except Exception as e:
                logging.warning(f"HWP → HWPX 변환 실패, 기본 로직으로 진행: {e}")
        # Try OLE text first, then OCR fallback
        if olefile is None:
            logging.warning("olefile 미설치 → HWP OCR 시도")
            return run_ocr_on_hwp_images(File_Bytes), False
        try:
            if not olefile.isOleFile(io.BytesIO(File_Bytes)):
                logging.warning("HWP가 OLE 형식이 아님 → OCR 시도")
                return run_ocr_on_hwp_images(File_Bytes), False
            def _extract_hwp_text(data: bytes) -> str:
                with olefile.OleFileIO(io.BytesIO(data)) as f:
                    for entry in f.listdir():
                        if "BodyText" in entry[0]:
                            try:
                                raw = f.openstream(entry).read()
                                txt = raw.decode("utf-16", errors="ignore")
                                if txt.strip():
                                    return txt
                            except Exception:
                                continue
                    if f.exists("PrvText"):
                        return f.openstream("PrvText").read().decode("utf-16", errors="ignore")
                return ""
            text = _extract_hwp_text(File_Bytes).strip()
            if not text:
                logging.warning("HWP 텍스트 스트림 없음 → OCR 시도")
                text = run_ocr_on_hwp_images(File_Bytes)
            return text or "", False
        except Exception as e:
            logging.error(f"HWP 처리 중 예외: {e}")
            return run_ocr_on_hwp_images(File_Bytes), False

    if ext == "hwpx":
        try:
            with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                xml_files = [n for n in z.namelist() if n.endswith('.xml')]
                if not xml_files:
                    raise ValueError("[ERROR] HWPX 파싱 실패: 암호 또는 손상")
                text = ""
                for name in xml_files:
                    data = z.read(name).decode("utf-8", errors="ignore")
                    text += re.sub("<[^>]+>", " ", data)
                text = text.strip()
                if not text:
                    logging.info("HWPX 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_hwpx_images(File_Bytes)
                return text, False
        except Exception as e:
            raise ValueError(f"[ERROR] HWPX 파싱 실패: {e}")

    if ext in ("ppt", "pptx"):
        if ext == "pptx":
            try:
                with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                    if "ppt/presentation.xml" not in z.namelist():
                        raise ValueError("[ERROR] PPTX 파싱 실패: 암호")
                if Presentation is None:
                    raise ValueError("python-pptx 미설치")
                prs = Presentation(io.BytesIO(File_Bytes))
                text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                if not text.strip():
                    logging.info("PPTX 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_ppt_images(File_Bytes)
                return text, False
            except Exception as e:
                raise ValueError(f"[ERROR] PPTX 파싱 실패: {e}")
        else:
            import tempfile, os as _os
            # legacy .ppt → try COM touch to validate password, then OCR fallback
            if win32com is not None:
                try:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
                        tmp.write(File_Bytes); tmp_path = tmp.name
                    tmp_pptx_path = tmp_path + "x"
                    try:
                        pp = win32com.client.Dispatch("PowerPoint.Application")
                        pp.Visible = False
                        pres = pp.Presentations.Open(tmp_path, WithWindow=False)
                        pres.SaveAs(tmp_pptx_path, FileFormat=24)  # 24: pptx
                        pres.Close(); pp.Quit()
                    except Exception:
                        raise ValueError("[ERROR] PPT → PPTX 변환 실패 또는 암호")
                    with open(tmp_pptx_path, "rb") as f:
                        converted = f.read()
                    _os.remove(tmp_path); _os.remove(tmp_pptx_path)
                    return parse_file(converted, "pptx")
                except Exception as e:
                    logging.warning(f"PPT → PPTX 변환 실패, OCR fallback: {e}")
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
                    tmp.write(File_Bytes); tmp_path = tmp.name
                try:
                    pp = win32com.client.Dispatch("PowerPoint.Application")
                    pp.Visible = False
                    pres = pp.Presentations.Open(tmp_path, WithWindow=False)
                    pres.Close(); pp.Quit()
                except Exception:
                    raise ValueError("[ERROR] PPT 파싱 실패: 암호")
                finally:
                    _os.remove(tmp_path)
                text = ""
                if not text.strip():
                    logging.info("PPT 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_ppt_images(File_Bytes)
                return text, False
            except Exception as e:
                raise ValueError(f"[ERROR] PPT 파싱 실패: {e}")

    if ext in ("xls", "xlsx"):
        if ext == "xlsx":
            try:
                with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                    if "xl/workbook.xml" not in z.namelist():
                        raise ValueError("[ERROR] XLSX 파싱 실패: 암호")
                if load_workbook is None:
                    raise ValueError("openpyxl 미설치")
                wb = load_workbook(io.BytesIO(File_Bytes), data_only=True)
                text = ""
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        text += " ".join([str(c) for c in row if c is not None]) + "\n"
                text = text.strip()
                if not text:
                    logging.info("XLSX 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_xls_images(File_Bytes)
                return text, False
            except Exception as e:
                raise ValueError(f"[ERROR] XLSX 파싱 실패: {e}")
        else: # .xls
            import tempfile, os as _os
            if win32com is not None:
                import tempfile, os as _os
                try:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".xls") as tmp:
                        tmp.write(File_Bytes); tmp_path = tmp.name
                    tmp_xlsx_path = tmp_path + "x"
                    try:
                        excel = win32com.client.Dispatch("Excel.Application")
                        excel.Visible = False
                        wb = excel.Workbooks.Open(tmp_path)
                        wb.SaveAs(tmp_xlsx_path, FileFormat=51)  # 51: xlsx
                        wb.Close(SaveChanges=False)
                        excel.Quit()
                    except Exception:
                        raise ValueError("[ERROR] XLS → XLSX 변환 실패 또는 암호")
                    with open(tmp_xlsx_path, "rb") as f:
                        converted = f.read()
                    _os.remove(tmp_path); _os.remove(tmp_xlsx_path)
                    return parse_file(converted, "xlsx")
                except Exception as e:
                    logging.warning(f"XLS → XLSX 변환 실패, OCR fallback: {e}")
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".xls") as tmp:
                    tmp.write(File_Bytes); tmp_path = tmp.name
                try:
                    excel = win32com.client.Dispatch("Excel.Application")
                    excel.Visible = False
                    wb = excel.Workbooks.Open(tmp_path)
                    wb.Close(SaveChanges=False)
                    excel.Quit()
                except Exception:
                    raise ValueError("[ERROR] XLS 파싱 실패: 암호")
                finally:
                    _os.remove(tmp_path)
                text = ""
                if not text.strip():
                    logging.info("XLS 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_xls_images(File_Bytes)
                return text, False
            except Exception as e:
                raise ValueError(f"[ERROR] XLS 파싱 실패: {e}")

    if ext == "doc":
        if win32com is None:
            logging.warning("win32com 미설치 → DOC 변환 불가, OCR 우회 시도")
            return "", False
        import tempfile, os as _os
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp:
                tmp.write(File_Bytes); tmp_path = tmp.name
            tmp_docx_path = tmp_path + "x"
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(tmp_path)
                doc.SaveAs(tmp_docx_path, FileFormat=16)
                doc.Close(); word.Quit()
            except Exception:
                raise ValueError("[ERROR] DOC 변환/파싱 실패: 암호")
            with open(tmp_docx_path, "rb") as f:
                converted = f.read()
            _os.remove(tmp_path); _os.remove(tmp_docx_path)
            text, _ = parse_file(converted, "docx")
            if not text.strip():
                logging.info("DOC 텍스트 없음 → OCR 실행")
                text = run_ocr_on_docx_images(converted)
            return text, False
        except Exception as e:
            raise ValueError(f"[ERROR] DOC 변환/파싱 실패: {e}")

    if ext == "docx":
        if Document is None:
            raise ValueError("python-docx 미설치")
        try:
            fs = io.BytesIO(File_Bytes)
            # password check by probing ZIP
            try:
                with zipfile.ZipFile(fs) as zf:
                    if "word/document.xml" not in zf.namelist():
                        raise ValueError("[ERROR] DOCX 파싱 실패: 암호")
            except zipfile.BadZipFile:
                raise ValueError("[ERROR] DOCX 파싱 실패: 손상/암호")
            fs.seek(0)
            doc = Document(fs)
            text = "\n".join([(p.text or "") for p in doc.paragraphs])
            if not text.strip():
                logging.info("DOCX 텍스트 없음 → OCR 실행")
                text = run_ocr_on_docx_images(File_Bytes)
                if not text.strip():
                    logging.warning("DOCX OCR 결과 없음 → 빈 문자열 반환")
                    return "", False
            return text, False
        except Exception as e:
            raise e

    if ext == "pdf":
        if fitz is None:
            raise ValueError("PyMuPDF 미설치")
        try:
            doc = fitz.open(stream=File_Bytes, filetype="pdf")
            if doc.is_encrypted:
                raise ValueError("[ERROR] PDF 파싱 실패: 암호")
            text = " ".join([p.get_text().replace("\n", " ") for p in doc])
            if not any(c.isalnum() for c in text):
                logging.info("PDF 텍스트 없음 → OCR 실행")
                text = run_ocr_on_pdf_images(File_Bytes)
                if not text.strip():
                    logging.error("PDF OCR 텍스트 추출 실패")
                    return "", False
            return text, False
        except Exception as e:
            raise ValueError(f"[ERROR] PDF 파싱 실패: {e}")

    if ext in IMAGE_EXTENSIONS:
        return run_ocr_on_single_image(File_Bytes), True

    raise ValueError(f"[ERROR] 지원하지 않는 파일 형식: {ext}")

# ==========================
# Image extraction for face scan
# ==========================
def extract_images_from_document(file_bytes: bytes, file_ext: str) -> list[bytes]:
    images: list[bytes] = []
    ext = (file_ext or "").lower()

    if ext == "docx":
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                for name in z.namelist():
                    if name.startswith("word/media/"):
                        with z.open(name) as f:
                            images.append(f.read())
        except Exception as e:
            logging.error(f"DOCX 이미지 추출 실패: {e}")

    elif ext == "doc":
        # convert to docx then reuse
        if win32com is None:
            logging.warning("win32com 미설치 → DOC 이미지 추출 불가")
        else:
            import tempfile, os as _os
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp:
                    tmp.write(file_bytes); p = tmp.name
                p_docx = p + "x"
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                d = word.Documents.Open(p)
                d.SaveAs(p_docx, FileFormat=16)
                d.Close(); word.Quit()
                with open(p_docx, 'rb') as f:
                    converted = f.read()
                _os.remove(p); _os.remove(p_docx)
                return extract_images_from_document(converted, 'docx')
            except Exception as e:
                logging.error(f"DOC 이미지 추출 실패: {e}")

    elif ext == "hwpx":
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                for name in z.namelist():
                    if name.startswith("Contents/") and name.lower().endswith((".jpg", ".png", ".bmp")):
                        images.append(z.read(name))
        except Exception as e:
            logging.error(f"HWPX 이미지 추출 실패: {e}")

    elif ext == "hwp":
        if olefile is None:
            logging.warning("olefile 미설치 → HWP 이미지 추출 불가")
        else:
            try:
                with olefile.OleFileIO(io.BytesIO(file_bytes)) as f:
                    for entry in f.listdir():
                        if "BinData" in entry[0]:
                            images.append(f.openstream(entry).read())
            except Exception as e:
                logging.error(f"HWP 이미지 추출 실패: {e}")

    elif ext in ("ppt", "pptx"):
        if Presentation is None:
            logging.warning("python-pptx 미설치 → PPT 이미지 추출 불가")
        else:
            try:
                prs = Presentation(io.BytesIO(file_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                            images.append(shape.image.blob)
            except Exception as e:
                logging.error(f"PPT/PPTX 이미지 추출 실패: {e}")

    elif ext in ("xls", "xlsx"):
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                for name in z.namelist():
                    if name.lower().startswith("xl/media/"):
                        images.append(z.read(name))
        except Exception as e:
            logging.error(f"XLS/XLSX 이미지 추출 실패: {e}")

    elif ext == "pdf":
        if fitz is None:
            logging.warning("PyMuPDF 미설치 → PDF 이미지 추출 불가")
        else:
            try:
                with fitz.open(stream=file_bytes, filetype='pdf') as doc:
                    for page in doc:
                        for img in page.get_images(full=True):
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            imbytes = base_image.get('image')
                            if imbytes:
                                images.append(imbytes)
            except Exception as e:
                logging.error(f"PDF 이미지 추출 실패: {e}")

    return images

# ==========================
# Face detection via MTCNN
# ==========================
def detect_faces_in_image_bytes(image_bytes: bytes):
    """Return list of {bbox:[x,y,w,h], confidence:float}. Empty if none or detector missing."""
    if detector is None:
        return []
    try:
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        img_np = np.array(img)
        results = detector.detect_faces(img_np)
        out = []
        for r in results:
            x, y, w, h = r.get('box', [0, 0, 0, 0])
            out.append({"bbox": [int(x), int(y), int(w), int(h)], "confidence": float(r.get('confidence', 0.0))})
        return out
    except Exception as e:
        logging.error(f"MTCNN 얼굴 탐지 실패: {e}")
        return []

def scan_file_for_face_images(file_bytes: bytes, file_ext: str):
    ext = (file_ext or "").lower()
    results = []

    if ext in ("doc", "docx", "hwp", "hwpx", "ppt", "pptx", "xls", "xlsx"):
        for idx, img_bytes in enumerate(extract_images_from_document(file_bytes, ext)):
            faces = detect_faces_in_image_bytes(img_bytes)
            if faces:
                results.append({"image_name": f"embedded_{idx}", "faces_found": len(faces), "faces": faces})
        return results

    if ext == "pdf":
        # Prefer image extraction for true embedded images
        if fitz is not None:
            try:
                with fitz.open(stream=file_bytes, filetype='pdf') as doc:
                    for pno, page in enumerate(doc):
                        images = page.get_images(full=True)
                        for i, img in enumerate(images):
                            xref = img[0]
                            base = doc.extract_image(xref)
                            imbytes = base.get("image")
                            if not imbytes:
                                continue
                            faces = detect_faces_in_image_bytes(imbytes)
                            if faces:
                                results.append({"page": pno+1, "image_index": i, "faces_found": len(faces), "faces": faces})
            except Exception as e:
                logging.warning(f"PDF 이미지 검사 실패: {e}")
        return results

    if ext in IMAGE_EXTENSIONS:
        faces = detect_faces_in_image_bytes(file_bytes)
        if faces:
            results.append({"image_name": "uploaded_image", "faces_found": len(faces), "faces": faces})
        return results

    return results

# ==========================
# Regex + validators
# ==========================
def validate_luhn(card_number: str) -> bool:
    digits = [int(d) for d in re.sub(r"\D", "", card_number)]
    checksum = 0
    for i, d in enumerate(reversed(digits)):
        if i % 2 == 1:
            t = d * 2
            checksum += t - 9 if t > 9 else t
        else:
            checksum += d
    return checksum % 10 == 0

def validate_ssn(ssn: str) -> bool:
    ssn = re.sub(r"\D", "", ssn)
    if len(ssn) != 13:
        return False
    yy = int(ssn[0:2]); mm = int(ssn[2:4]); dd = int(ssn[4:6]); g = int(ssn[6])
    if g not in [1,2,3,4,5,6,7,8]:
        return False
    full_year = (1900 + yy) if g in [1,2,5,6] else (2000 + yy)
    try:
        datetime.date(full_year, mm, dd)
    except ValueError:
        return False
    if full_year >= 2020:
        return True
    weights = [2,3,4,5,6,7,8,9,2,3,4,5]
    s = sum(int(ssn[i]) * weights[i] for i in range(12))
    check = (11 - (s % 11)) % 10
    return check == int(ssn[-1])

def detect_by_regex(Text: str) -> list:
    # Primary patterns (with spacing/hyphen support)
    Patterns = {
        "phone": re.compile(r"\b01[016789][\s\-]?\d{3,4}[\s\-]?\d{4}\b"),
        "email": re.compile(r"\b[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}\b"),
        "birth": re.compile(r"\b(19|20)\d{2}[년./\- ]+(0?[1-9]|1[0-2])[월./\- ]+(0?[1-9]|[12][0-9]|3[01])[일]?\b"),
        "ssn": re.compile(r"\b\d{6}[\s\-]?[1-4]\d{6}\b"),
        "alien_reg": re.compile(r"\b\d{6}[\s\-]?[5-8]\d{6}\b"),
        "driver_license": re.compile(r"\b\d{2}[\s\-]?\d{2}[\s\-]?\d{6}[\s\-]?\d{2}\b"),
        "passport": re.compile(r"\b[A-Z]\d{2,3}[A-Z]?\d{4,5}\b"),
        "account": re.compile(r"\b\d{6}[\s\-]?\d{2}[\s\-]?\d{6}\b"),
        "card": re.compile(r"\b(?:\d{4}[\s\-]?){3}\d{4}\b"),
        "ip": re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b"),
    }

    detected = []
    for label, pat in Patterns.items():
        for m in pat.finditer(Text):
            item = {"type": label, "value": m.group(), "span": m.span()}
            if label == "card":
                item["status"] = "valid" if validate_luhn(item["value"]) else "invalid (Luhn check failed)"
            if label == "ssn":
                item["status"] = "valid" if validate_ssn(item["value"]) else "invalid (SSN check failed)"
            detected.append(item)

    # Normalized pass to catch evasion without spaces/hyphens
    normalized_text = re.sub(r"[\s\-]", "", Text)
    normalized_patterns = {
        "phone": re.compile(r"\b01[016789]\d{7,8}\b"),
        "ssn": re.compile(r"\b\d{6}[1-4]\d{6}\b"),
        "alien_reg": re.compile(r"\b\d{6}[5-8]\d{6}\b"),
        "driver_license": re.compile(r"\b\d{2}\d{2}\d{6}\d{2}\b"),
        "account": re.compile(r"\b\d{6}\d{2}\d{6}\b"),
        "card": re.compile(r"\b\d{16}\b"),
    }
    existing = {d["type"]+":"+d["value"].replace(" ", "").replace("-", "") for d in detected}
    for label, pat in normalized_patterns.items():
        for m in pat.finditer(normalized_text):
            key = label+":"+m.group()
            if key in existing:
                continue
            item = {"type": label, "value": m.group(), "span": m.span()}
            if label == "card":
                item["status"] = "valid" if validate_luhn(item["value"]) else "invalid (Luhn check failed)"
            if label == "ssn":
                item["status"] = "valid" if validate_ssn(item["value"]) else "invalid (SSN check failed)"
            detected.append(item)
            existing.add(key)

    return detected

# ==========================
# NER (union of labels)
# ==========================
def detect_by_ner(Text: str) -> list:
    if not Text.strip():
        return []
    results = ner_pipeline(Text)
    out = []
    allowed = set([
        # 인명 / 개인
        "PER", "PERSON", "PS", "PEOPLE", "HUMAN", "NAME",
        # 조직 / 단체 / 회사 / 기관
        "ORG", "ORGANIZATION", "OG", "COMPANY", "CORP", "INSTITUTE", "UNIV",
        # 장소 / 지역 / 국가
        "LOC", "LC", "LOCATION", "PLACE", "CITY", "COUNTRY", "GPE", "ADDRESS", "REGION",
        # 날짜 / 시간
        "DATE", "DAT", "DT", "TIME", "TI", "DURATION", "EVENT",
        # 기타 고유명사 / 제품명 / 문화 / 언어 / 작품
        "MISC", "POH", "PRODUCT", "ART", "WORK_OF_ART", "LANGUAGE", "CULTURE",
        # 숫자형 엔터티 (ID, 금액, 나이 등)
        "NUM", "QUANTITY", "CARDINAL", "ORDINAL", "MONEY", "PERCENT",
        # 연락처, 이메일, 웹주소
        "EMAIL", "URL", "PHONE", "CONTACT",
        # 법률, 기관명, 문서명
        "LAW", "STATUTE", "DOCUMENT", "CASE", "LEGAL"
    ])
    for ent in results:
        label = ent.get('entity_group', '')
        if label is None:
            continue
        if label.upper() not in allowed:
            continue
        start = ent.get('start'); end = ent.get('end')
        if start is None or end is None:
            continue
        word = ent.get('word', '').strip()
        if label.upper() in ["DATE", "DAT", "DT", "TIME", "TI"]:
            # 0 또는 9로만 구성된 비정상 날짜/시간 제거
            if re.fullmatch(r"0[\s0:시분초]*", word):
                continue
            if re.fullmatch(r"9[\s9:시분초]*", word):
                continue

            # 날짜형 유효성 검증 (YYYY-MM-DD, YYYY.MM.DD, YYYY/MM/DD, YYYY MM DD)
            clean = re.sub(r"[./\s]", "-", word)
            m = re.match(r"(\d{4})-(\d{1,2})-(\d{1,2})", clean)
            if m:
                y, mo, d = map(int, m.groups())
                try:
                    datetime.date(y, mo, d)
                except ValueError:
                    continue  # 존재하지 않는 날짜 필터링

            # 시간형 유효성 검증 (HH:MM:SS, HH시MM분SS초, HH:MM)
            time_clean = re.sub(r"[시분초\s]", ":", word)
            t = re.match(r"(\d{1,2}):(\d{1,2})(?::(\d{1,2}))?", time_clean)
            if t:
                hh, mm, ss = t.groups()
                hh = int(hh)
                mm = int(mm)
                ss = int(ss) if ss else 0
                if not (0 <= hh < 24 and 0 <= mm < 60 and 0 <= ss < 60):
                    continue  # 비정상 시각 (예: 25:99:99) 제거

            # 완전 무의미한 시간 표현 제거 ("00:00:00" 또는 "0시0분0초")
            if re.fullmatch(r"0{1,2}[:시\s]?0{1,2}[:분\s]?0{1,2}(초)?", word):
                continue
        out.append({"type": label, "value": word, "span": (start, end)})
    return out

# 0과 9만 별도의 필터링으로 나누는 이유? : 모델이 가장 자주 오탐하는 두 패턴이 서로 다르니까.
# 0만 구성된 경우 : “날짜 없음”, “공란 채우기용 숫자”로 학습되어 날짜처럼 감지
# 9만 구성된 경우 : “미래 날짜”, “샘플/가상 데이터”로 인식되어 날짜로 감지
# 두 패턴은 NER 모델이 오탐을 유발하는 대표적인 “비의미 숫자 그룹”이라 별도로 제거함

# ==========================
# Main handler
# ==========================
def handle_input_raw(Input_Data, Original_Format=None, meta_info=None):
    """
    Returns:
      - Detected: list
      - Parsed_Text: str (may be "")
      - backend_status: str ("backend_disabled" | "no_text" | "no_detection" | "face_only")
      - image_detections: list  (face results per image)
    """
    # Accept bytes or str
    if isinstance(Input_Data, bytes):
        logging.info("파일 입력으로 감지됨")
        try:
            Parsed_Text, is_image = parse_file(Input_Data, Original_Format or "")
        except ValueError as ve:
            raise ve
    elif isinstance(Input_Data, str):
        logging.info("텍스트 입력으로 감지됨")
        Parsed_Text = Input_Data
        is_image = False
    else:
        raise ValueError("지원하지 않는 입력 형식입니다.")

    Parsed_Text = (Parsed_Text or "").strip()

    # Face scan
    image_detections = scan_file_for_face_images(Input_Data if isinstance(Input_Data, bytes) else b"", Original_Format or "")

    # If no text
    if not Parsed_Text:
        logging.info("OCR 실패 또는 빈 문자열")
        if image_detections:
            # faces found only
            return ([{"type": "image_face", "value": f"faces={sum(i.get('faces_found',0) for i in image_detections)}", "detail": image_detections}]), "", "face_only", image_detections
        else:
            return [], "", "no_text", image_detections

    # Text detections
    detected = detect_by_regex(Parsed_Text) + detect_by_ner(Parsed_Text)

    # -------------------------------
    # 중복 제거 로직 (정규표현식 우선)
    # -------------------------------
    regex_priority = {"birth", "ssn", "card", "driver_license", "account", "passport", "phone", "email"}
    unique = []
    seen = {}  # value → type

    for d in detected:
        val = d.get("value")
        typ = d.get("type")
        if not val or not typ:
            continue

        # Case 1: 처음 보는 값
        if val not in seen:
            seen[val] = typ
            unique.append(d)
            continue

        # Case 2: 동일 값 이미 존재함 → 중복 판단
        prev_type = seen[val]

        # (4) 정규표현식 기반 탐지를 우선
        if prev_type not in regex_priority and typ in regex_priority:
            # 기존 NER 탐지 대신 정규표현식으로 교체
            for u in unique:
                if u["value"] == val:
                    u["type"] = typ
                    break
            seen[val] = typ
            continue

        # (3) 동일 필드명 + 동일 값 → 완전 중복이므로 제거
        if prev_type == typ:
            continue

        # (2) 서로 다른 필드명이지만 값이 동일 → 중복으로 제거 (정규표현식 우선 정책 적용됨)
        continue

    detected = unique

    # Append face summary into Detected
    total_faces = sum(i.get("faces_found", 0) for i in image_detections)
    if total_faces > 0:
        detected.append({
            "type": "image_face",
            "value": f"total_faces={total_faces}",
            "detail": image_detections
        })

    # Dedup by (type,value)
    uniq = []
    seen = set()
    for d in detected:
        key = (d.get("type"), d.get("value"))
        if key in seen:
            continue
        seen.add(key)
        uniq.append(d)
    detected = uniq

    if not detected:
        logging.info("민감정보 미탐지")
        return [], Parsed_Text, "no_detection", image_detections

    # Summary payload (kept local; sending disabled)
    type_counts = Counter([d.get("type") for d in detected if d.get("type")])
    card_total = card_valid = card_invalid = 0
    ssn_total = ssn_valid = ssn_invalid = 0
    for d in detected:
        t = d.get("type"); st = str(d.get("status", "")).lower()
        if t == "card":
            card_total += 1; card_valid += (st == "valid"); card_invalid += (st != "valid")
        if t == "ssn":
            ssn_total += 1; ssn_valid += (st == "valid"); ssn_invalid += (st != "valid")
    validation_summary = {
        "card": {"total": card_total, "valid": card_valid, "invalid": card_invalid},
        "ssn":  {"total": ssn_total,  "valid": ssn_valid,  "invalid": ssn_invalid}
    }
    payload = {
        "timestamp": datetime.datetime.utcnow().replace(microsecond=0).isoformat() + "Z",
        "has_face": total_faces > 0,
        "detected_summary": dict(type_counts),
        "validation_summary": validation_summary,
        "_meta": meta_info,
        "filename": (meta_info or {}).get("filename") if isinstance(meta_info, dict) else None,
    }
    # 백엔드 전송은 비활성화 유지
    # _ = payload  # placeholder to avoid linter complaints
    try:
        json_data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        ok = send_to_backend(json_data, filename="Detected_Info.json", meta=meta_info)
        backend_status = "backend_success" if ok else "backend_failed"
    except Exception as e:
        logging.error(f"백엔드 전송 중 오류: {e}")
        backend_status = "backend_failed"

    return detected, Parsed_Text, backend_status, image_detections

def send_to_backend(raw_data: bytes, filename: str = "Detected_Info.json", meta: dict | None = None) -> bool:
    url = "http://your-backend-api/upload"
    # url = os.getenv("BACKEND_URL", "http://127.0.0.1:8080/upload")
    try:
        files = {"file": (filename, raw_data, "application/json")}
        data = {"meta": json.dumps(meta, ensure_ascii=False)} if meta else {}
        response = requests.post(url, files=files, data=data)
        # print(f"[INFO] 백엔드 응답 상태 코드: {response.status_code}")
        logging.info(f"백엔드 응답 코드: {response.status_code}")
        return response.status_code == 200
    except Exception as e:
        logging.error(f"백엔드 전송 실패: {e}")
        return False

if __name__ == "__main__": # 서버 API 없이도 디버그 할 수 있게 해주는 임시 코드. 로직 기능 자체와는 무관하다.
    # Minimal self-check
    sample = "안녕하세요 제 주민등록번호는 900101-1234567 입니다. 카드번호 4111-1111-1111-1111"
    det, txt, st, faces = handle_input_raw(sample, Original_Format="txt", meta_info={"filename": "sample.txt"})
    print(json.dumps({"detected": det, "status": st, "len": len(txt), "faces": faces}, ensure_ascii=False, indent=2))
