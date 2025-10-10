import io
import re
import requests
import logging
import numpy as np
import fitz
import easyocr
import zipfile
import datetime
import cv2
import os
import face_recognition
import json
import olefile
import win32com.client
import tempfile
from PIL import Image
from transformers import AutoTokenizer, AutoModelForTokenClassification, pipeline
from docx import Document
from collections import Counter
from pptx import Presentation
from openpyxl import load_workbook

logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(message)s')
# NER_MODEL_NAME = "Leo97/KoELECTRA-small-v3-modu-ner"
NER_MODEL_NAME = "soddokayo/klue-roberta-base-ner"
# NER_MODEL_NAME = "amoeba04/koelectra-small-v3-privacy-ner"
# NER_MODEL_NAME = "Davlan/xlm-roberta-base-ner-hrl"
os.environ["HF_HUB_DOWNLOAD_TIMEOUT"] = "180"
os.environ["HF_HUB_ETAG_TIMEOUT"] = "180"
ner_tokenizer = AutoTokenizer.from_pretrained(NER_MODEL_NAME)
ner_model = AutoModelForTokenClassification.from_pretrained(NER_MODEL_NAME)
ner_pipeline = pipeline("ner", model=ner_model, tokenizer=ner_tokenizer, grouped_entities=True)
reader = easyocr.Reader(['ko', 'en'])
IMAGE_EXTENSIONS = ("png", "jpg", "jpeg", "bmp", "webp", "tiff", "tif")

def parse_file(File_Bytes: bytes, File_Ext: str) -> str:
    File_Ext = File_Ext.lower()
    if File_Ext == "txt":
        try:
            return File_Bytes.decode("utf-8"), None
        except UnicodeDecodeError:
            return File_Bytes.decode("cp949"), None
        
    elif File_Ext == "hwp":
        try:
            # olefile이 아니거나 스트림 읽기 실패 시 암호/손상으로 간주
            if not olefile.isOleFile(io.BytesIO(File_Bytes)):
                raise ValueError("[ERROR] HWP 파싱 실패: 파일이 OLE 형식이 아닙니다.")
            with olefile.OleFileIO(io.BytesIO(File_Bytes)) as f:
                try:
                    encoded_text = f.openstream("PrvText").read()
                except Exception as e:
                    raise ValueError("[ERROR] HWP 파싱 실패: 암호로 보호된 문서입니다.")
            text = encoded_text.decode("utf-16")
            if not text.strip():
                print("[INFO] HWP 텍스트 없음 → OCR 실행")
                text = run_ocr_on_hwp_images(File_Bytes)
            return text, None
        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"[ERROR] HWP 파싱 실패: {e}")
        
    elif File_Ext == "hwpx":
        try:
            with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                # hwpx는 XML이 있어야 정상. 없으면 암호 또는 손상으로 간주
                xml_files = [n for n in z.namelist() if n.endswith(".xml")]
                if not xml_files:
                    raise ValueError("[ERROR] HWPX 파싱 실패: 암호로 보호된 문서입니다.")
                text = ""
                for name in xml_files:
                    data = z.read(name).decode("utf-8", errors="ignore")
                    text += re.sub("<[^>]+>", " ", data)
                if not text.strip():
                    print("[INFO] HWPX 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_hwpx_images(File_Bytes)
                return text, None
        except zipfile.BadZipFile:
            raise ValueError("[ERROR] HWPX 파싱 실패: 파일이 손상되었거나 암호로 보호된 문서입니다.")
        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"[ERROR] HWPX 파싱 실패: {e}")
        
    elif File_Ext in ["ppt", "pptx"]:
        if File_Ext == "pptx":
            try:
                with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                    if "ppt/presentation.xml" not in z.namelist():
                        raise ValueError("[ERROR] PPTX 파싱 실패: 암호로 보호된 문서입니다.")
                prs = Presentation(io.BytesIO(File_Bytes))
                text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                if not text.strip():
                    print("[INFO] PPTX 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_ppt_images(File_Bytes)
                return text, None
            except zipfile.BadZipFile:
                raise ValueError("[ERROR] PPTX 파싱 실패: 파일이 손상되었거나 암호로 보호된 문서입니다.")
            except ValueError:
                raise
            except Exception as e:
                raise ValueError(f"[ERROR] PPTX 파싱 실패: {e}")
        else:  # legacy .ppt (binary)
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
                    tmp.write(File_Bytes); tmp_path = tmp.name
                try:
                    # PowerPoint COM으로 열기 시 비밀번호가 있으면 예외 발생
                    pp = win32com.client.Dispatch("PowerPoint.Application")
                    pp.Visible = False
                    pres = pp.Presentations.Open(tmp_path, WithWindow=False)
                    pres.Close()
                    pp.Quit()
                except Exception as e:
                    raise ValueError("[ERROR] PPT 파싱 실패: 암호로 보호된 문서입니다.")
                finally:
                    os.remove(tmp_path)
                text = ""
                if not text.strip():
                    print("[INFO] PPT 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_ppt_images(File_Bytes)
                return text, None
            except ValueError:
                raise
            except Exception as e:
                raise ValueError(f"[ERROR] PPT 파싱 실패: {e}")
        
    elif File_Ext in ["xls", "xlsx"]:
        if File_Ext == "xlsx":
            try:
                with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                    if "xl/workbook.xml" not in z.namelist():
                        raise ValueError("[ERROR] XLSX 파싱 실패: 암호로 보호된 문서입니다.")
                wb = load_workbook(io.BytesIO(File_Bytes), data_only=True)
                text = ""
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        text += " ".join([str(cell) for cell in row if cell]) + "\n"
                text = text.strip()
                if not text:
                    print("[INFO] XLSX 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_xls_images(File_Bytes)
                return text, None
            except zipfile.BadZipFile:
                raise ValueError("[ERROR] XLSX 파싱 실패: 파일이 손상되었거나 암호로 보호된 문서입니다.")
            except ValueError:
                raise
            except Exception as e:
                raise ValueError(f"[ERROR] XLSX 파싱 실패: {e}")
        else:  # legacy .xls
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".xls") as tmp:
                    tmp.write(File_Bytes); tmp_path = tmp.name
                try:
                    excel = win32com.client.Dispatch("Excel.Application")
                    excel.Visible = False
                    wb = excel.Workbooks.Open(tmp_path)
                    wb.Close(SaveChanges=False)
                    excel.Quit()
                except Exception:
                    raise ValueError("[ERROR] XLS 파싱 실패: 암호로 보호된 문서입니다.")
                finally:
                    os.remove(tmp_path)
                    text = ""
                    if not text.strip():
                        print("[INFO] XLS 텍스트 없음 → OCR 실행")
                        text = run_ocr_on_xls_images(File_Bytes)
                    return text, None
            except ValueError:
                raise
            except Exception as e:
                raise ValueError(f"[ERROR] XLS 파싱 실패: {e}")
        
    elif File_Ext == "doc":
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp_doc:
                tmp_doc.write(File_Bytes)
                tmp_doc_path = tmp_doc.name
            tmp_docx_path = tmp_doc_path + "x"
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(tmp_doc_path)
                doc.SaveAs(tmp_docx_path, FileFormat=16)
                doc.Close()
                word.Quit()
            except Exception:  # COM 열기 실패 → 암호문서로 간주
                raise ValueError("[ERROR] DOC 변환/파싱 실패: 암호로 보호된 문서입니다.")

            with open(tmp_docx_path, "rb") as f:
                converted_bytes = f.read()
            os.remove(tmp_doc_path)
            os.remove(tmp_docx_path)

            # --- OCR fallback 추가 부분 ---
            text, _ = parse_file(converted_bytes, "docx")
            if not text.strip():
                print("[INFO] DOC 텍스트 없음 → OCR 실행")
                text = run_ocr_on_doc_images(File_Bytes)
            return text, None
        # -----------------------------

        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"[ERROR] DOC 변환/파싱 실패: {e}")
    
    elif File_Ext == "docx":
        try:
            file_stream = io.BytesIO(File_Bytes)
            try:
                with zipfile.ZipFile(file_stream) as zf:
                    if "word/document.xml" not in zf.namelist():
                        raise ValueError("[ERROR] DOCX 파싱 실패: 암호로 보호된 문서입니다.")
            except zipfile.BadZipFile:
                raise ValueError("[ERROR] DOCX 파싱 실패: 파일이 손상되었거나 암호로 보호된 문서입니다.")
            except RuntimeError as e:
                if "password" in str(e).lower():
                    raise ValueError("[ERROR] DOCX 파싱 실패: 암호로 보호된 문서입니다.")
                else:
                    raise
            file_stream.seek(0)
            doc = Document(file_stream)
            text = "\n".join([para.text or "" for para in doc.paragraphs])
            if not text.strip():
                print("[INFO] 텍스트 없음 → OCR 실행")
                text = run_ocr_on_docx_images(File_Bytes)
                if not text.strip():
                    print("[WARN] OCR 결과 없음 → 빈 문자열 반환 (얼굴 탐지 여부는 별도 처리)")
                    return "", None   # 여기서 예외를 던지지 않음
            return text, None
        except Exception as e:
            # raise ValueError("f[ERROR] DOCX 파싱 실패: {e}")
            raise e

    elif File_Ext == "pdf":
        try:
            doc = fitz.open(stream=File_Bytes, filetype="pdf")
            if doc.is_encrypted:
                raise ValueError("[ERROR] PDF 파싱 실패 : 암호로 보호된 PDF 문서입니다.")
            text = " ".join([page.get_text().replace("\n", " ") for page in doc])
            if not any(char.isalnum() for char in text):
                print("[INFO] 텍스트 없음 → OCR 실행")
                text = run_ocr_on_pdf_images(File_Bytes)
                if not text.strip():
                    print("[ERROR] OCR 텍스트 추출 실패, 빈 문자열 반환")
                    return "", None
            return text, None
        except Exception as e:
            if any(keyword in str(e).lower() for keyword in ["암호", "cannot open broken document", "not a pdf", "password"]):
                raise ValueError(str(e))
            raise ValueError("[ERROR] PDF 파싱 실패")
    
    elif File_Ext in ["png", "jpg", "jpeg", "bmp", "webp", "tif", "tiff"]:
        try:
            image = Image.open(io.BytesIO(File_Bytes))
            result = reader.readtext(np.array(image))
            ocr_text = " ".join([box[1] for box in result])
            return ocr_text, None
        except Exception as e:
            raise ValueError(f"[ERROR] 이미지 OCR 실패: {e}")
        
def detect_faces(image_bytes: bytes, save_dir="processed_faces", filename_prefix="face_detected"):
    try:
        os.makedirs(save_dir, exist_ok=True)
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        np_img = np.array(image)
        face_locations = face_recognition.face_locations(np_img)
        print(f"[DEBUG] face_locations={face_locations}", flush=True)

        if len(face_locations) > 0:
            # for (top, right, bottom, left) in face_locations:
            #     face_region = np_img[top:bottom, left:right]
            #     small = cv2.resize(face_region, (10, 10), interpolation=cv2.INTER_LINEAR)
            #     mosaic = cv2.resize(small, (right - left, bottom - top), interpolation=cv2.INTER_NEAREST)
            #     np_img[top:bottom, left:right] = mosaic
            # -> 모자이크 로직, 그냥 주석 처리.
            # filename = f"{filename_prefix}_{datetime.datetime.now(datetime.timezone.utc).isoformat()}.jpg"
            timestamp = datetime.datetime.utcnow().replace(microsecond=0).isoformat() + "Z"
            filename = f"{filename_prefix}_{timestamp.replace(':', '-')}.jpg"
            save_path = os.path.join(save_dir, filename)
            cv2.imwrite(save_path, cv2.cvtColor(np_img, cv2.COLOR_RGB2BGR))
            return True, save_path
        return False, None

    except Exception as e:
        print(f"[ERROR] 얼굴 탐지 실패 (face_recognition): {e}")
        return False, None
    
def extract_images_from_document(file_bytes: bytes, file_ext: str) -> list:
    images = []
    file_ext = file_ext.lower()

    if file_ext == "docx":
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as docx_zip:
                image_files = [f for f in docx_zip.namelist() if f.startswith("word/media/")]
                for image_name in image_files:
                    with docx_zip.open(image_name) as image_file:
                        images.append(image_file.read())
        except Exception as e:
            print(f"[ERROR] DOCX 이미지 추출 실패: {e}")

    elif file_ext == "doc":
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            tmp_docx_path = tmp_path + "x"
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(tmp_path)
            doc.SaveAs(tmp_docx_path, FileFormat=16)
            doc.Close()
            word.Quit()
            with open(tmp_docx_path, "rb") as f:
                converted_bytes = f.read()
            os.remove(tmp_path); os.remove(tmp_docx_path)
            return extract_images_from_document(converted_bytes, "docx")
        except Exception as e:
            print(f"[ERROR] DOC 이미지 추출 실패: {e}")

    elif file_ext == "hwpx":
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                for name in z.namelist():
                    if name.startswith("Contents/") and name.lower().endswith((".jpg", ".png", ".bmp")):
                        images.append(z.read(name))
        except Exception as e:
            print(f"[ERROR] HWPX 이미지 추출 실패: {e}")

    elif file_ext == "hwp":
        try:
            with olefile.OleFileIO(io.BytesIO(file_bytes)) as f:
                for entry in f.listdir():
                    if "BinData" in entry[0]:
                        data = f.openstream(entry).read()
                        images.append(data)
        except Exception as e:
            print(f"[ERROR] HWP 이미지 추출 실패: {e}")

    elif file_ext in ["ppt", "pptx"]:
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # PICTURE
                        image = shape.image
                        images.append(image.blob)
        except Exception as e:
            print(f"[ERROR] PPT/PPTX 이미지 추출 실패: {e}")

    elif file_ext in ["xls", "xlsx"]:
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                for name in z.namelist():
                    if name.lower().startswith("xl/media/"):
                        images.append(z.read(name))
        except Exception as e:
            print(f"[ERROR] XLS/XLSX 이미지 추출 실패: {e}")

    elif file_ext == "pdf":
        try:
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                for page in doc:
                    pix = page.get_pixmap()
                    img_bytes = pix.tobytes("ppm")
                    images.append(img_bytes)
        except Exception as e:
            print(f"[ERROR] PDF 이미지 추출 실패: {e}")

    return images

def run_ocr_on_doc_images(file_bytes: bytes) -> str:
    """DOC 파일을 DOCX로 변환 후 OCR"""
    ocr_text = ""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        tmp_docx_path = tmp_path + "x"
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(tmp_path)
        doc.SaveAs(tmp_docx_path, FileFormat=16)
        doc.Close()
        word.Quit()
        with open(tmp_docx_path, "rb") as f:
            converted_bytes = f.read()
        os.remove(tmp_path); os.remove(tmp_docx_path)
        ocr_text = run_ocr_on_docx_images(converted_bytes)
    except Exception as e:
        print(f"[ERROR] DOC 이미지 OCR 실패: {e}")
    return ocr_text.strip()

def run_ocr_on_docx_images(file_bytes):
    ocr_text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as docx_zip:
            image_files = [f for f in docx_zip.namelist() if f.startswith("word/media/")]
            for image_name in image_files:
                with docx_zip.open(image_name) as image_file:
                    image = Image.open(image_file)
                    result = reader.readtext(np.array(image))
                    for box in result:
                        ocr_text += box[1] + "\n"
    except Exception as e:
        print(f"[ERROR] DOCX 이미지 OCR 실패: {e}")
    return ocr_text.replace("\n", " ").replace(":", "").replace(",", "").strip()

def run_ocr_on_hwp_images(file_bytes: bytes) -> str:
    """HWP 문서 내 BinData 이미지에 대해 OCR 수행"""
    ocr_text = ""
    try:
        with olefile.OleFileIO(io.BytesIO(file_bytes)) as f:
            for entry in f.listdir():
                if "BinData" in entry[0]:
                    data = f.openstream(entry).read()
                    try:
                        img = Image.open(io.BytesIO(data))
                        result = reader.readtext(np.array(img))
                        for box in result:
                            ocr_text += box[1] + "\n"
                    except Exception:
                        continue
    except Exception as e:
        print(f"[ERROR] HWP 이미지 OCR 실패: {e}")
    return ocr_text.replace("\n", " ").strip()

def run_ocr_on_hwpx_images(file_bytes: bytes) -> str:
    """HWPX 파일 내 이미지 OCR"""
    ocr_text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            for name in z.namelist():
                if name.startswith("Contents/") and name.lower().endswith((".jpg", ".png", ".bmp")):
                    img = Image.open(io.BytesIO(z.read(name)))
                    result = reader.readtext(np.array(img))
                    for box in result:
                        ocr_text += box[1] + "\n"
    except Exception as e:
        print(f"[ERROR] HWPX 이미지 OCR 실패: {e}")
    return ocr_text.replace("\n", " ").strip()

def run_ocr_on_ppt_images(file_bytes: bytes) -> str:
    """PPT/PPTX 슬라이드 내 이미지 OCR"""
    ocr_text = ""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    try:
                        img = Image.open(io.BytesIO(shape.image.blob))
                        result = reader.readtext(np.array(img))
                        for box in result:
                            ocr_text += box[1] + "\n"
                    except Exception:
                        continue
    except Exception as e:
        print(f"[ERROR] PPT/PPTX 이미지 OCR 실패: {e}")
    return ocr_text.replace("\n", " ").strip()

def run_ocr_on_xls_images(file_bytes: bytes) -> str:
    """XLS/XLSX 내 이미지 OCR"""
    ocr_text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            for name in z.namelist():
                if name.lower().startswith("xl/media/"):
                    img = Image.open(io.BytesIO(z.read(name)))
                    result = reader.readtext(np.array(img))
                    for box in result:
                        ocr_text += box[1] + "\n"
    except Exception as e:
        print(f"[ERROR] XLS/XLSX 이미지 OCR 실패: {e}")
    return ocr_text.replace("\n", " ").strip()

def run_ocr_on_pdf_images(pdf_bytes: bytes) -> str:
    ocr_text = ""
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        for idx, page in enumerate(doc):
            pix = page.get_pixmap()
            img = Image.open(io.BytesIO(pix.tobytes("ppm")))
            try:
                result = reader.readtext(np.array(img))
                for box in result:
                    ocr_text += box[1] + "\n"
            except Exception as e:
                print(f"[ERROR] EasyOCR 실패 (페이지 {idx}): {e}")
    return ocr_text.replace("\n", " ").replace(":", "").replace(",", "").strip()

def validate_luhn(card_number: str) -> bool:
    digits = [int(d) for d in re.sub(r"\D", "", card_number)]
    checksum = 0
    reverse = digits[::-1]
    for i, digit in enumerate(reverse):
        if i % 2 == 1:
            doubled = digit * 2
            checksum += doubled - 9 if doubled > 9 else doubled
        else:
            checksum += digit
    return checksum % 10 == 0

def validate_ssn(ssn: str) -> bool:
    ssn = re.sub(r"\D", "", ssn)
    if len(ssn) != 13:
        return False

    yy = int(ssn[0:2])
    mm = int(ssn[2:4])
    dd = int(ssn[4:6])
    gender_digit = int(ssn[6])

    # 성별 코드 유효성
    if gender_digit not in [1,2,3,4,5,6,7,8]:
        return False

    # 세기 구분
    if gender_digit in [1,2,5,6]:
        full_year = 1900 + yy
    elif gender_digit in [3,4,7,8]:
        full_year = 2000 + yy
    else:
        return False

    # 날짜 유효성
    try:
        datetime.date(full_year, mm, dd)
    except ValueError:
        return False

    # 2020년 이후 출생자는 체크섬 패스
    if full_year >= 2020:
        return True

    # 체크섬 검증 (2020년 이전 출생자만)
    weights = [2,3,4,5,6,7,8,9,2,3,4,5]
    sum_val = sum(int(ssn[i]) * weights[i] for i in range(12))
    check_digit = (11 - (sum_val % 11)) % 10
    return check_digit == int(ssn[-1])

def detect_by_regex(Text: str) -> list:
    Patterns = {
        "phone": re.compile(r"\b01[016789][-\s]?\d{3,4}[-\s]?\d{4}\b"),
        "email": re.compile(r"\b[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}\b"),
        "birth": re.compile(r"\b(19|20)\d{2}[년./\- ]+(0?[1-9]|1[0-2])[월./\- ]+(0?[1-9]|[12][0-9]|3[01])[일]?\b"),
        "ssn": re.compile(r"\b\d{6}[-\s]?[1-4]\d{6}\b"),
        "alien_reg": re.compile(r"\b\d{6}[-\s]?[5-8]\d{6}\b"),
        "driver_license": re.compile(r"\b\d{2}[-\s]?\d{2}[-\s]?\d{6}[-\s]?\d{2}\b"),  # [-\s]? = 하이픈/띄어쓰기 선택적
        "passport": re.compile(r"\b[A-Z]\d{2,3}[A-Z]?\d{4,5}\b"),
        "account": re.compile(r"\b\d{6}[-\s]?\d{2}[-\s]?\d{6}\b"),
        "card": re.compile(r"\b(?:\d{4}[-\s]?){3}\d{4}\b"),
        "ip": re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b")
    }

    detected = []
    for p_type, pattern in Patterns.items():
        for match in pattern.finditer(Text):
            raw_value = match.group()
            start, end = match.span()
            result = {
                "type": p_type,
                "value": raw_value,
                "span": (start, end)
            }
            if p_type == "card":
                if not validate_luhn(raw_value):
                    result["status"] = "invalid (Luhn check failed)"
                else:
                    result["status"] = "valid"
            elif p_type == "ssn":
                if not validate_ssn(raw_value):
                    result["status"] = "invalid (SSN check failed)"
                else:
                    result["status"] = "valid"
            detected.append(result)
    return detected

def detect_by_ner(Text: str) -> list:
    Results = ner_pipeline(Text)
    detected = []
    for Entity in Results:
        Label = Entity['entity_group']
        Word = Entity['word']
        Start = Entity.get('start')
        End = Entity.get('end')
        if Start is None or End is None:
            continue
        if Label.upper() in [
            # 인명 / 개인
            "PER", "PERSON", "PS", "PEOPLE", "HUMAN", "NAME",
            # 조직 / 단체 / 회사 / 기관
            "ORG", "ORGANIZATION", "OG", "COMPANY", "CORP", "INSTITUTE", "UNIV",
            # 장소 / 지역 / 국가
            "LOC", "LC", "LOCATION", "PLACE", "CITY", "COUNTRY", "GPE", "ADDRESS", "REGION",
            # 날짜 / 시간
            "DATE", "DAT", "DT", "TIME", "TI", "DURATION", "EVENT",
            # 기타 고유명사 / 제품명 / 문화 / 언어 / 작품
            "MISC", "POH", "PRODUCT", "ART", "WORK_OF_ART", "LANGUAGE", "CULTURE",
            # 숫자형 엔터티 (ID, 금액, 나이 등)
            "NUM", "QUANTITY", "CARDINAL", "ORDINAL", "MONEY", "PERCENT",
            # 연락처, 이메일, 웹주소 (일부 영어 모델 탐지)
            "EMAIL", "URL", "PHONE", "CONTACT",
            # 법률, 기관명, 문서명 (법률문서 기반 NER용)
            "LAW", "STATUTE", "DOCUMENT", "CASE", "LEGAL"
            ]:
            detected.append({"type": Label, "value": Word, "span": (Start, End)})
    return detected

# def apply_masking(original_text: str, detected: list) -> str:
#     masked_text = list(original_text)
#     length = len(masked_text)

#     for item in detected:
#         start, end = item.get("span", (None, None))
#         if start is None or end is None:
#             continue
#         if not isinstance(start, int) or not isinstance(end, int):
#             continue
#         if start < 0 or end > length or start >= end:
#             continue
#         if end > len(masked_text):
#             continue
#         try:
#             masked_text[start:end] = "*" * (end - start)
#         except IndexError as e:
#             print(f"[WARNING] 마스킹 실패 (index error): {e} | span=({start}, {end})")
#             continue

#     return ''.join(masked_text)
# -> 마스킹 로직, 그냥 주석 처리

def send_to_backend(encrypted_data: bytes, filename: str = "Detected_Info", meta: dict | None = None) -> bool:
    url = "http://your-backend-api/upload"
    try:
        files = {"File": (filename, encrypted_data)}
        data = {}
        if meta:
            data["meta"] = json.dumps(meta, ensure_ascii=False)
        # response = requests.post(url, files={"File": (filename, encrypted_data)})
        response = requests.post(url, files=files, data=data)
        print(f"[INFO] 백엔드 응답 상태 코드: {response.status_code}")
        return response.status_code == 200
    except requests.exceptions.RequestException as e:
        print(f"[ERROR] 백엔드 전송 실패: {e}")
        return False
    
def handle_input_raw(Input_Data, Original_Format=None, meta_info=None):
    if isinstance(Input_Data, bytes):
        print("[INFO] 파일 입력으로 감지됨")
        Parsed_Text, _ = parse_file(Input_Data, Original_Format)
    elif isinstance(Input_Data, str):
        print("[INFO] 텍스트 입력으로 감지됨")
        Parsed_Text = Input_Data
    else:
        raise ValueError("지원하지 않는 입력 형식입니다.")
    Parsed_Text = Parsed_Text.strip()
    face_detected = False
    face_path = None
    if Original_Format in IMAGE_EXTENSIONS:
        face_detected, face_path = detect_faces(Input_Data)
        if face_detected: # 얼굴 탐지 결과를 Detected에 추가
            return [{"type": "face", "value": "image"}], "", "face_only", face_path
    elif Original_Format in ["pdf", "docx", "hwp", "hwpx", "ppt", "pptx", "xls", "xlsx", "doc"]:
         # 내부 이미지 추출 후 하나씩 face 분석
        for img in extract_images_from_document(Input_Data, Original_Format):
            face_detected, face_path = detect_faces(img)
            if face_detected:
                break
    if not Parsed_Text.strip():
        print("[INFO] OCR 실패 또는 빈 문자열 → 전송 생략")
        if face_detected:
            return [{"type": "face", "value": "image"}], "", "face_only", face_path
        else:
            return [], "", "no_text", None  
    Detected = detect_by_regex(Parsed_Text) + detect_by_ner(Parsed_Text)
    unique_detected = []
    seen_values = set()
    for d in Detected:
        value = d.get("value")
        if value not in seen_values:
            seen_values.add(value)
            unique_detected.append(d)
    Detected = unique_detected
    if face_detected and face_path:
        Detected.insert(0, {"type": "face", "value": "image"})
    if not Detected:
        print("[INFO] 민감정보 미탐지 → 백엔드 전송 생략")
        # return Detected, "", False, face_path
        return [], "", "no_detection", face_path
    type_counts = Counter([d.get("type") for d in Detected if d.get("type")])
    card_total = card_valid = card_invalid = 0
    ssn_total = ssn_valid = ssn_invalid = 0
    for d in Detected:
        t = d.get("type")
        st = d.get("status", "").lower()
        if t == "card":
            card_total += 1
            if st == "valid":
                card_valid += 1
            else:
                card_invalid += 1
        elif t == "ssn":
            ssn_total += 1
            if st == "valid":
                ssn_valid += 1
            else:
                ssn_invalid += 1
    validation_summary = {
        "card": {"total": card_total, "valid": card_valid, "invalid": card_invalid},
        "ssn":  {"total": ssn_total,  "valid": ssn_valid,  "invalid": ssn_invalid}
        }
    payload = {
        "timestamp": datetime.datetime.utcnow().replace(microsecond=0).isoformat() + "Z",
        # "source_format": Original_Format,
        "has_face": face_detected,
        "detected_summary": dict(type_counts), # "detected": [], <- 원래 이걸 쓰던 상황
        "validation_summary": validation_summary,
        "_meta": meta_info
        # "detected": [
        #     {"type": d.get("type"), "value": d.get("value"), "span": d.get("span")}
        #     for d in Detected
        # ]
    }
    payload["filename"] = meta_info.get("filename") if meta_info else None
    # for d in Detected:
    #     if d.get("type") != "face":  # 얼굴 항목은 제외
    #         payload["detected"].append({
    #             "type": d.get("type"),
    #             "value": d.get("value"),
    #             "span": d.get("span")
    #             })
    ok = send_to_backend(json.dumps(payload, ensure_ascii=False).encode("utf-8"), filename="Detected_Items.json")
    return Detected, Parsed_Text, "sent_ok" if ok else "send_fail", face_path
    # encrypted = encrypt_data(Parsed_Text.encode("utf-8"), Key)
    # backend_status = send_to_backend(encrypted, filename="Detected_Info")
    # return Detected, Parsed_Text, backend_status, face_path