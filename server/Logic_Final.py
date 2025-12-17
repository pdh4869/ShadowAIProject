# =============================
# File: Logic_Updated_Final.py
# Desc: Logic_Final.py 버전에 다음 기능들을 통합한 최종 개선 버전입니다.
#       1. 조합 위험도 분석 강화: 'critical' 등급 추가 및 규칙 세분화
#       2. 이름 탐지 정확도 향상: NER 결과의 신뢰도(confidence)를 기반으로 오탐지 필터링
#       3. 구버전 파일 지원 확대: win32com을 이용한 .doc 파일 파싱 기능 추가
#       4. (신규) GIF 파일 내 텍스트 탐지(OCR) 기능 추가
# =============================

import io
import re
import os
import logging
import zipfile
import tempfile
import datetime
import numpy as np
from collections import Counter
from concurrent.futures import ThreadPoolExecutor

# --- 필수 라이브러리 ---
from transformers import AutoTokenizer, AutoModelForTokenClassification, pipeline

# --- 선택 라이브러리 (설치되지 않아도 기본 기능 동작) ---
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None
try:
    import easyocr
except ImportError:
    easyocr = None
try:
    from PIL import Image, ImageSequence
except ImportError:
    Image = None
    ImageSequence = None
try:
    from docx import Document
except ImportError:
    Document = None
try:
    from mtcnn import MTCNN
except ImportError:
    MTCNN = None
try:
    import olefile
except ImportError:
    olefile = None
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    import win32com.client
except ImportError:
    win32com = None
try:
    import xlrd # .xls 지원을 위해 추가
except ImportError:
    xlrd = None

# 로깅
DEBUG_MODE = os.getenv("PII_DEBUG", "false").lower() == "true"
logging.basicConfig(
    level=logging.DEBUG if DEBUG_MODE else logging.INFO,
    format='[%(levelname)s] %(message)s'
)

# ==========================
# NER 모델 로딩
# ==========================
HF_TOKEN = os.getenv("HF_TOKEN", None)
NER_MODEL_NAME = "soddokayo/klue-roberta-base-ner"
print(f"[INFO] NER 모델 로딩 중: {NER_MODEL_NAME}")
if HF_TOKEN:
    ner_tokenizer = AutoTokenizer.from_pretrained(NER_MODEL_NAME, token=HF_TOKEN)
    ner_model = AutoModelForTokenClassification.from_pretrained(NER_MODEL_NAME, token=HF_TOKEN)
    print("[INFO] [OK] 허깅페이스 토큰 인증 완료")
else:
    print("[WARN] HF_TOKEN 환경 변수 없음 - 공개 모델로 시도")
    ner_tokenizer = AutoTokenizer.from_pretrained(NER_MODEL_NAME)
    ner_model = AutoModelForTokenClassification.from_pretrained(NER_MODEL_NAME)

ner_pipeline = pipeline("ner", model=ner_model, tokenizer=ner_tokenizer, grouped_entities=True)
print("[INFO] [OK] NER 모델 로딩 완료")

# EasyOCR 초기화
reader = None
if easyocr and Image is not None:
    try:
        reader = easyocr.Reader(['ko', 'en'], gpu=False)
        print("[INFO] EasyOCR 초기화 완료")
    except Exception as e:
        print(f"[WARN] EasyOCR 초기화 실패: {e}")
else:
    print("[WARN] EasyOCR 또는 PIL 미설치")

# MTCNN 초기화
detector = None
if MTCNN:
    try:
        detector = MTCNN()
        print("[INFO] MTCNN 초기화 완료")
    except Exception as e:
        print(f"[WARN] MTCNN 초기화 실패: {e}")
else:
    print("[WARN] mtcnn 라이브러리가 설치되지 않았습니다.")

# ==========================
# 검증 함수 (Luhn/주민등록)
# ==========================

def validate_luhn(card_number: str) -> bool:
    digits = [int(d) for d in re.sub(r"\D", "", card_number)]
    if len(digits) < 13:
        return False
    checksum = 0
    for i, d in enumerate(reversed(digits)):
        if i % 2 == 1:
            t = d * 2
            checksum += t - 9 if t > 9 else t
        else:
            checksum += d
    return checksum % 10 == 0

def validate_ssn(ssn: str) -> bool:
    ssn = re.sub(r"\D", "", ssn)
    if len(ssn) != 13:
            return False  # 주민등록번호는 13자리여야 합니다.
    yy, mm, dd, g = int(ssn[0:2]), int(ssn[2:4]), int(ssn[4:6]), int(ssn[6])
    if g not in [1,2,3,4,5,6,7,8]:
        return False
    century = 1900 if g in [1,2,5,6] else 2000
    try:
        datetime.date(century + yy, mm, dd)
    except ValueError:
        return False
    # 2020년 이후 출생자 일부 예외 허용
    if (century + yy) >= 2020:
        return True
    weights = [2,3,4,5,6,7,8,9,2,3,4,5]
    s = sum(int(ssn[i]) * weights[i] for i in range(12))
    check_digit = (11 - (s % 11)) % 10
    return check_digit == int(ssn[-1])

# ==========================
# 정규식 패턴
# ==========================

COMPILED_PATTERNS = {
    "phone": re.compile(r'\b(?:010[\s-]?\d{3,4}[\s-]?\d{4}|0(?:2|3[1-3]|4[1-4]|5[1-5]|6[1-4]|70)[\s-]?\d{3,4}[\s-]?\d{4})\b'),
    "email": re.compile(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'),
    "birth": re.compile(r"(?<!\d)(19[0-9]{2}|20[0-2][0-9])[년./\-\s]+(0?[1-9]|1[0-2])[월./\-\s]+(0?[1-9]|[12][0-9]|3[01])[일]?(?!\d)"),
    "ssn": re.compile(r"(?<!\d)\d{6}[\s\-]?[1-4]\d{6}(?!\d)"),
    "alien_reg": re.compile(r"(?<!\d)\d{6}[\s\-]?[5-8]\d{6}(?!\d)"),
    "driver_license": re.compile(r"(?<!\d)(1[1-9]|2[0-8])[\s\-]?\d{2}[\s\-]?\d{6}[\s\-]?\d{2}(?!\d)"),
    "passport": re.compile(r"\b[A-Z]\d{2,3}[A-Z]?\d{4,5}\b"),
    "account": re.compile(r"(?<!\d)\d{6}[\s\-]?\d{2}[\s\-]?\d{6}(?!\d)"),
    "card": re.compile(r"(?<!\d)(?:\d{4}[\s\-]?){3}\d{4}(?!\d)"),
    "ip": re.compile(r"(?<!\d)(?:\d{1,3}\.){3}\d{1,3}(?!\d)")
}
COMPILED_NORMALIZED_PATTERNS = {
    "phone_normalized": re.compile(r"(?<!\d)(?:010\d{7,8}|02\d{7,8}|0(?:3[1-3]|4[1-4]|5[1-5]|6[1-4]|70)\d{7,8})(?!\d)"),
    "ssn_normalized": re.compile(r"(?<!\d)\d{6}[1-4]\d{6}(?!\d)"),
    "alien_reg_normalized": re.compile(r"(?<!\d)\d{6}[5-8]\d{6}(?!\d)"),
    "driver_license_normalized": re.compile(r"(?<!\d)(1[1-9]|2[0-8])\d{10}(?!\d)"),
    "account_normalized": re.compile(r"(?<!\d)\d{14}(?!\d)"),
    "card_normalized": re.compile(r"(?<!\d)\d{16}(?!\d)")
}

KOREAN_SURNAMES = {'김','이','박','최','정','강','조','윤','장','임','한','오','서','신','권','황','안','송','류','전','홍','고','문','양','손','배','백','허','남','심','노','하','곽','성','차','주','우','구','라','진','유'}
NAME_WHITELIST = {'홍길동', '유재석'}

# 조직명 화이트리스트 (NER이 놓치는 특정 회사명)
ORG_WHITELIST = {'홈플러스', '협진축산', '홈플러스간석점'}

# ==========================
# 파일명 마스킹
# ==========================

def mask_pii_in_filename(filename: str) -> tuple:
    detected_items = detect_by_regex(filename)
    detected_items.extend(detect_by_ner(filename))

    if not detected_items:
        return (filename, None)

    unique_items = []
    detected_items.sort(key=lambda x: x.get('span', (9999, 9999))[0])
    last_end = -1
    for item in detected_items:
        span = item.get('span')
        if not span:
            continue
        s, e = span
        if s >= last_end:
            unique_items.append(item)
            last_end = e

    if not unique_items:
        return (filename, None)

    type_map = {
        'ssn': '주민등록번호','email':'이메일','phone':'전화번호','card':'카드번호',
        'driver_license':'운전면허','account':'계좌번호','passport':'여권번호','alien_reg':'외국인등록번호'
    }
    detected_types = set()
    for it in unique_items:
        t = it.get('type','')
        display_type = '이름' if t in ['PS','PER','NER_PS','NER_PER'] else type_map.get(t, t)
        detected_types.add(display_type)

    masked = filename
    unique_items.sort(key=lambda x: x.get('span', (0,0))[1], reverse=True)
    for it in unique_items:
        s, e = it.get('span', (None,None))
        if s is None:
            continue
        val = it.get('value','')
        masked = masked[:s] + ('*' * len(val)) + masked[e:]

    return (masked, ', '.join(sorted(detected_types)) if detected_types else None)

# ==========================
# 정규식 / NER / 준식별자
# ==========================

def detect_by_regex(Text: str) -> list:
    normalized_text = re.sub(r'[\s\-]', '', Text)
    detected = []
    seen_values = set()  # 중복 방지
    for label, pattern in COMPILED_PATTERNS.items():
        for match in pattern.finditer(Text):
            # 전화번호 검증
            if label == "phone":
                start, end = match.span()
                matched_text = match.group()
                phone_digits = re.sub(r'\D', '', matched_text)
                
                # 전화번호 길이 검증 (9~11자리)
                if len(phone_digits) < 9 or len(phone_digits) > 11:
                    continue
                
                # 유효한 전화번호 형식인지 먼저 검증
                valid_prefixes = ['010', '02', '031', '032', '033', '041', '042', '043', '044', 
                                  '051', '052', '053', '054', '055', '061', '062', '063', '064', '070']
                if not any(phone_digits.startswith(prefix) for prefix in valid_prefixes):
                    continue
                
                # 앞뒤 좁은 컨텍스트 확인
                context_start = max(0, start - 10)
                context_end = min(len(Text), end + 10)
                context = Text[context_start:context_end]
                
                # 날짜 패턴이 바로 인접해 있으면 제외
                if re.search(r'(19|20)\d{2}[년월일\-/\.]', context):
                    continue
                
                # 앞뒤에 영문자가 붙어있으면 제외
                if (start > 0 and Text[start-1].isalpha()) or (end < len(Text) and Text[end].isalpha()):
                    continue
                
                # 연속된 숫자 체크
                if len(phone_digits) > 1:
                    is_sequential = all(int(phone_digits[i]) == (int(phone_digits[i-1]) + 1) % 10 for i in range(1, len(phone_digits)))
                    if is_sequential:
                        continue
            
            # 생년월일 검증: 키워드 기반 필터링
            if label == "birth":
                start, end = match.span()
                # 앞뒤 50글자 범위에서 키워드 찾기
                context_start = max(0, start - 50)
                context_end = min(len(Text), end + 50)
                context = Text[context_start:context_end].lower()
                
                # 생년월일 관련 키워드
                birth_keywords = ['생년월일', '생일', '출생', '생년', 'birth', 'dob', 'date of birth']
                # 제외 키워드 (입사일 등)
                exclude_keywords = ['입사', '퇴사', '계약', '신고', '등록', '수정', '발급', '승인', '승인일', '가입', '신청', 'join', 'hire', 'contract', 'register']
                
                has_birth_keyword = any(kw in context for kw in birth_keywords)
                has_exclude_keyword = any(kw in context for kw in exclude_keywords)
                
                # 생년월일 키워드가 있고 제외 키워드가 없을 때만 탐지
                if not has_birth_keyword or has_exclude_keyword:
                    continue
            
            matched_value = match.group()
            # 중복 체크
            value_key = f"{label}:{matched_value}"
            if value_key in seen_values:
                continue
            seen_values.add(value_key)
            
            item = {"type": label, "value": matched_value, "span": match.span()}
            if label == "card":
                item["status"] = "valid" if validate_luhn(item["value"]) else "invalid (Luhn)"
            if label == "ssn":
                item["status"] = "valid" if validate_ssn(item["value"]) else "invalid (SSN)"
            detected.append(item)

    existing = set()
    for d in detected:
        normalized_val = re.sub(r'[\s-]', '', d["value"])
        existing.add(normalized_val)
        if d["type"] == "phone" and normalized_val.startswith('+'):
            existing.add(normalized_val[1:])

    for label, pattern in COMPILED_NORMALIZED_PATTERNS.items():
        original = label.replace("_normalized", "")
        for m in pattern.finditer(normalized_text):
            nv = m.group()
            if nv in existing:
                continue
            if original == "phone" and not nv.startswith('+') and f"+{nv}" in existing:
                continue
            
            # 전화번호 normalized 패턴 추가 검증
            if original == "phone":
                # 유효한 전화번호 형식인지 검증
                valid_prefixes = ['010', '02', '031', '032', '033', '041', '042', '043', '044', 
                                  '051', '052', '053', '054', '055', '061', '062', '063', '064', '070']
                if not any(nv.startswith(prefix) for prefix in valid_prefixes):
                    continue

            vp = re.compile(r'[\s-]*'.join(list(nv)))
            rm = vp.search(Text)
            if rm:
                original_value = rm.group()
                if original == "card" and re.search(r'\d{4}[\s-]\d{2}[\s-]\d{2}', original_value):
                    continue

                # 중복 체크
                value_key = f"{original}:{original_value}"
                if value_key not in seen_values:
                    seen_values.add(value_key)
                    item = {"type": original, "value": original_value, "span": rm.span()}
                    if original == "card":
                        item["status"] = "valid" if validate_luhn(item["value"]) else "invalid (Luhn)"
                    if original == "ssn":
                        item["status"] = "valid" if validate_ssn(item["value"]) else "invalid (SSN)"
                    detected.append(item)
                existing.add(nv)
            else:
                value_key = f"{original}:{nv}"
                if value_key not in seen_values:
                    seen_values.add(value_key)
                    detected.append({"type": original, "value": nv, "span": None})
                existing.add(nv)
    return detected

def detect_by_ner(Text: str) -> list:
    if not Text.strip():
        return []
    
    Detected = []
    detected_names = set()
    detected_orgs = set()  # 조직명 중복 방지
    
    # 정규식 보완 탐지 비활성화 (NER만 사용)
    
    # 제외할 단어 (헤더, 지명, 일반 명사)
    exclude_words = {
        # 헤더/라벨
        '성명', '주소', '이름', '성함', '직위', '직급', '부서', '소속',
        # 지명 (구/동/로)
        '강남구', '서초구', '송파구', '강동구', '강서구', '양천구', '구로구', '영등포구', '동작구', '관악구', '서대문구', '마포구', '용산구', '성동구', '광진구', '동대문구', '중랑구', '성북구', '강북구', '도봉구', '노원구', '은평구', '종로구', '중구',
        '한강대로', '테헤란로', '강남대로', '논현로', '봉은사로', '선릉로', '역삼로', '언주로', '도산대로', '압구정로', '서초대로', '반포대로', '사평대로', '효령로', '방배로', '동작대로', '상도로', '노량진로', '여의대로', '국회대로', '의사당대로', '마포대로', '서강대로', '독막로', '월드컵로', '성산로', '진흥로', '장한로',
        # 직위
        '사원', '대리', '과장', '차장', '부장', '이사', '상무', '전무', '부사장', '사장', '주임', '선임', '책임', '수석', '팀장', '실장', '본부장',
        # 일반 명사 (오탐지 방지)
        '오늘', '내일', '어제', '정보', '성격', '장점', '주요', '진료', '최적화', '이해할', '이야기를', '하였으며', '하안동', '홍콩'
    }
    
    # 정규식으로 완전한 주소 패턴 탐지
    address_pattern = re.compile(r'(서울특별시|부산광역시|대구광역시|인천광역시|광주광역시|대전광역시|울산광역시|세종특별자치시|경기도|강원도|충청북도|충청남도|전라북도|전라남도|경상북도|경상남도|제주특별자치도)\s*[가-힣]+(?:시|군|구)\s*[가-힣A-Za-z0-9]+(?:로|길)\s*[0-9]+(?:,\s*[A-Za-z가-힣]+)?(?:,\s*[A-Za-z0-9가-힣]+(?:층|호|동))?')
    detected_addresses = set()
    for match in address_pattern.finditer(Text):
        addr = match.group().strip().rstrip(',')
        if addr and addr not in detected_addresses:
            Detected.append({"type": "LC", "value": addr, "span": match.span()})
            detected_addresses.add(addr)
    
    # 화이트리스트 조직명 탐지 (NER 보완)
    for org_name in ORG_WHITELIST:
        if org_name in Text and org_name not in detected_orgs:
            start_idx = Text.find(org_name)
            Detected.append({"type": "ORG", "value": org_name, "span": (start_idx, start_idx + len(org_name))})
            detected_orgs.add(org_name)
    
    try:
        # NER 모델 512 토큰 제한 해결: 분할 처리
        ner_results = []
        max_length = 500  # 안전 마진
        if len(Text) > max_length:
            chunks = [Text[i:i+max_length] for i in range(0, len(Text), max_length)]
            for chunk in chunks:
                ner_results.extend(ner_pipeline(chunk))
        else:
            ner_results = ner_pipeline(Text)
        
        for whitelist_name in NAME_WHITELIST:
            if whitelist_name in Text:
                start_idx = Text.find(whitelist_name)
                Detected.append({"type": "PS", "value": whitelist_name, "span": (start_idx, start_idx + len(whitelist_name))})
                # 화이트리스트로 추가한 이름은 중복 방지를 위해 detected_names 집합에 추가
                detected_names.add(whitelist_name.replace(" ", ""))

        for entity in ner_results:
            Label = entity['entity_group'].upper()
            # 정규화: PER/PERSON 등 인명 레이블을 'PS'로 통일하여 동일 인물의 중복 탐지를 방지
            if Label in ('PER', 'PERSON'):
                Label = 'PS'
            Word = entity['word'].replace('##', '')
            Start = entity.get('start')
            End = entity.get('end')

            if Start is None or End is None:
                continue
            
            # LC(주소)는 정규식으로 이미 처리했으므로 스킵
            if Label == 'LC':
                continue

            if Label in ['PS', 'PER']:
                clean_word = Word.replace(" ", "").strip()
                if len(clean_word) <= 1 or clean_word.isdigit() or not any('\uac00' <= c <= '\ud7a3' for c in clean_word):
                    continue
                
                # 숫자 포함 제외
                if any(c.isdigit() for c in clean_word):
                    continue
                
                # 제외 단어 필터 (정규식과 동일)
                if clean_word in exclude_words:
                    continue
                
                # 직위 키워드 제외
                position_keywords = ['사원', '대리', '과장', '차장', '부장', '이사', '상무', '전무', '부사장', '사장', '주임', '선임', '책임', '수석', '부수석', '원장', '부원장', '국장', '부국장', '실장', '팀장']
                if clean_word in position_keywords:
                    continue
                
                org_keywords = ['회사', '전자', '그룹', '기업', '주식회사', '(주)', '㉼', '학교', '대학교', '대학', '고등학교', '중학교', '초등학교', '병원', '의원', '센터', '연구소', '재단', '협회', '은행', '부서', '팀', '본부', '지점', '영업소', '축산', '농장', '목장', '마트', '플러스', '점포', '상회']
                if any(kw in clean_word for kw in org_keywords):
                    if clean_word not in detected_orgs:
                        Detected.append({"type": "ORG", "value": Word, "span": (Start, End)})
                        detected_orgs.add(clean_word)
                    continue
                
                # 3글자 한글 이름만 허용
                if len(clean_word) != 3:
                    continue
                
                # 성씨로 시작하지 않으면 제외
                if clean_word[0] not in KOREAN_SURNAMES:
                    continue
                
                # 일반 명사 패턴 제외 (조사 포함)
                if clean_word.endswith(('을', '를', '이', '가', '은', '는', '에', '의', '로', '와', '과')):
                    continue
                
                # 중복 방지
                if clean_word in detected_names:
                    continue
                
                # 필터 통과한 이름만 추가
                Detected.append({"type": Label, "value": Word, "span": (Start, End)})
                detected_names.add(clean_word)
                continue
            
            if Label in ['ORG', 'OG']:
                clean_org = Word.replace(" ", "").strip()
                
                if '○' in clean_org or '□' in clean_org or clean_org in ['고등학교', '중학교', '초등학교', '대학교', '대학원', '대학']:
                    continue
                
                # 오탐지 필터: "-" 또는 숫자로 시작하는 조직명 제외
                if clean_org.startswith('-') or clean_org.startswith('- ') or (clean_org and clean_org[0].isdigit()):
                    continue
                
                # Word 원본도 체크 (공백 포함)
                if Word.strip().startswith('-') or Word.strip().startswith('- '):
                    continue
                
                # 중복 체크 (부분 문자열 포함)
                if clean_org in detected_orgs:
                    continue
                if any(clean_org in existing or existing in clean_org for existing in detected_orgs):
                    continue
                
                if len(clean_org) >= 2:
                    org_split_keywords = ['팀', '부', '부서', '본부', '지점', '센터', '연구소']
                    split_orgs = []
                    # '-' 포함된 조직명은 split 하지 않음
                    if ' ' in Word and '-' not in Word:
                        parts = Word.split(' ', 1)
                        if len(parts) == 2:
                            split_orgs = [parts[0].strip(), parts[1].strip()]
                    if not split_orgs:
                        for keyword in org_split_keywords:
                            if keyword in Word and not Word.endswith(keyword):
                                idx = Word.find(keyword)
                                if idx > 0:
                                    part1 = Word[:idx].strip()
                                    part2 = Word[idx:].strip()
                                    if part1 and part2:
                                        split_orgs = [part1, part2]
                                        break
                    if split_orgs:
                        for org in split_orgs:
                            clean_split = org.replace(" ", "").strip()
                            if len(clean_split) >= 2 and clean_split not in detected_orgs:
                                Detected.append({"type": "ORG", "value": org, "span": (Start, End)})
                                detected_orgs.add(clean_split)
                    else:
                        Detected.append({"type": Label, "value": Word, "span": (Start, End)})
                        detected_orgs.add(clean_org)
                continue
            
            if Label == 'LOC':
                Detected.append({"type": Label, "value": Word, "span": (Start, End)})

    except Exception as e:
        logging.warning(f"NER 파이프라인 오류: {e}")
            
    return Detected

def detect_quasi_identifiers(text: str) -> list:
    detected = []
    # 직책 탐지
    position_keywords = ['사원', '대리', '과장', '차장', '부장', '이사', '상무', '전무', '부사장', '사장', '주임', '선임', '책임', '수석', '부수석', '원장', '부원장', '국장', '부국장', '실장', '팀장', '본부장']
    position_pattern = re.compile(r'\b(' + '|'.join(position_keywords) + r')\b')
    for match in position_pattern.finditer(text):
        detected.append({"type":"position","value":match.group(),"span":match.span()})
    return detected

# ==========================
# 조합 위험도 (상세 메시지 버전)
# ==========================

def categorize_detection(item_type: str) -> str:
    # image_face는 민감정보(sensitive)로 유지
    if item_type in ['image_face']:
        return 'sensitive'
    # 식별자(identifier): 금융/민감한 직접 식별자들만 포함
    if item_type in ['phone', 'email', 'ssn', 'alien_reg', 'driver_license', 'passport', 'card', 'account']:
        return 'identifier'
    # 준식별자(quasi): 조직/주소/직위 등 기존 준식별자에 "이름(PS)"과 "ip"를 포함하도록 변경
    # 학번(student_id)은 준식별자 목록에서 제외
    if item_type in ['ORG','OG','birth','LC','position', 'LOC', 'PS', 'PER', 'ip']:
        return 'quasi'
    return 'other'

def _translate_type(item_type: str) -> str:
    type_map = {
        'PS': '이름', 'PER': '이름', 'image_face': '얼굴',
        'phone': '전화번호', 'email': '이메일', 'ssn': '주민번호',
        'alien_reg': '외국인번호', 'driver_license': '면허번호',
        'passport': '여권번호', 'card': '카드번호', 'account': '계좌번호', 'ip': 'IP주소',
        'ORG': '조직명', 'OG': '조직명', 'student_id': '학번', 'birth': '생년월일',
        'LC': '주소', 'LOC': '주소', 'position': '직위'
    }
    return type_map.get(item_type, item_type)

def analyze_combination_risk(detected_items, text):
    # 변경된 정책:
    # - 개인 식별 의심(combination risk)은 오직 '준식별자'들의 조합에서만 판단합니다.
    # - 식별자(identifier) 항목은 조합 판단에서 제외합니다.
    # - 준식별자에는 'PS'(이름)와 'ip'를 포함하며, 'student_id'는 준식별자에서 제외합니다.
    # - 같은 종류의 준식별자만 여러개 존재하는 경우(예: 이름 2개)에는 조합위험으로 간주하지 않습니다.

    if not detected_items or len(detected_items) < 2:
        return None

    # 그룹화: 우리는 준식별자(quasi)만 사용
    quasis = [it for it in detected_items if categorize_detection(it.get('type','')) == 'quasi']
    q_cnt = len(quasis)
    if q_cnt < 2:
        return None

    # 집계: 타입별 개수와 고유 타입 수
    q_types = [q.get('type') for q in quasis]
    unique_q_types = list(dict.fromkeys(q_types))
    unique_count = len(unique_q_types)

    # 같은 종류의 준식별자만 탐지된 경우(예: 이름만 2개)는 조합위험으로 간주하지 않음
    if unique_count <= 1:
        return None

    # 이제 준식별자 중 서로 다른 종류가 2종 이상 존재하면 조합위험으로 판단
    # 등급은 단순화: 서로다른 준식별자 2종 이상이면 'high'
    risk_level = 'high'
    # 메시지 구성 (기존 스타일 유지)
    q_types_trans = sorted(list(set([_translate_type(t) for t in unique_q_types])))
    q_str = f"준식별자({','.join(q_types_trans)}){q_cnt}건"
    risk_msg = f"{q_str} → 준식별자 조합으로 개인 특정 가능성" if q_cnt > 0 else None
    risk_items = quasis

    return {
        'level': risk_level,
        'message': risk_msg,
        'items': risk_items,
        'counts': {'quasi': q_cnt}
    }


# ==========================
# OCR
# ==========================

def run_ocr_on_single_image(image_bytes: bytes) -> str:
    if reader is None or Image is None:
        return ""
    try:
        # BytesIO 객체 생성 및 위치 초기화
        img_io = io.BytesIO(image_bytes)
        img_io.seek(0)
        
        # 이미지 열기 시도
        img = Image.open(img_io)
        
        # 이미지가 유효한지 확인
        img.verify()
        
        # verify() 후에는 다시 열어야 함
        img_io.seek(0)
        img = Image.open(img_io)
        
        # RGB로 변환
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        result = reader.readtext(np.array(img))
        return "\n".join([b[1] for b in result]).strip()
    except Exception:
        return ""


def _ocr_task(args):
    return run_ocr_on_single_image(args[1])


def run_ocr_on_docx_images(file_bytes):
    if reader is None or Image is None:
        return ""
    try:
        from PIL import ImageEnhance
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            image_files = [n for n in z.namelist() if n.startswith("word/media/")]
            if not image_files:
                return ""
            ocr_text = ""
            for image_name in image_files:
                try:
                    img = Image.open(io.BytesIO(z.read(image_name))).convert('RGB')
                    # 이미지 전처리: 대비 증가
                    enhancer = ImageEnhance.Contrast(img)
                    img = enhancer.enhance(2.0)
                    # OCR 실행 (신뢰도 임계값 낮춤)
                    result = reader.readtext(np.array(img), detail=1, paragraph=False)
                    for box in result:
                        if box[2] > 0.1:  # 신뢰도 10% 이상
                            ocr_text += box[1] + "\n"
                except Exception:
                    continue
            return ocr_text.strip()
    except Exception as e:
        print(f"[ERROR] DOCX 이미지 OCR 실패: {e}")
        return ""


def _process_pdf_image_ocr(args):
    pno, img_bytes = args
    return run_ocr_on_single_image(img_bytes)

def run_ocr_on_pdf_images(pdf_bytes: bytes) -> str:
    if reader is None or fitz is None:
        return ""
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        # 중복 이미지 제거 (xref 기반)
        seen_xrefs = set()
        image_tasks = []
        
        for pno, page in enumerate(doc):
            for img in page.get_images(full=True):
                xref = img[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                
                base_image = doc.extract_image(xref)
                img_bytes = base_image.get("image")
                
                # 작은 이미지 스킵 (5KB 이상만)
                if img_bytes and len(img_bytes) > 5000:
                    image_tasks.append((pno, img_bytes))
        
        # 병렬 처리
        if image_tasks:
            max_workers = min(8, os.cpu_count() or 4)
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                texts = list(executor.map(_process_pdf_image_ocr, image_tasks))
            return "\n".join(texts)
        return ""
    except Exception as e:
        print(f"[ERROR] PDF 이미지 OCR 실패: {e}")
        return ""


def run_ocr_on_hwp_images(hwp_bytes: bytes) -> str:
    if reader is None or olefile is None:
        return ""
    try:
        ole = olefile.OleFileIO(io.BytesIO(hwp_bytes))
        tasks = [(e, ole.openstream(e).read()) for e in ole.listdir() if e[0] == "BinData"]
        if not tasks:
            return ""
        with ThreadPoolExecutor() as ex:
            return "\n".join(ex.map(_ocr_task, tasks))
    except Exception as e:
        print(f"[ERROR] HWP 이미지 OCR 실패: {e}")
        return ""


def run_ocr_on_pptx_images(pptx_bytes: bytes) -> str:
    if reader is None:
        return ""
    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
            tasks = [(n, z.read(n)) for n in z.namelist() if n.startswith("ppt/media/")]
            if not tasks:
                return ""
            with ThreadPoolExecutor() as ex:
                return "\n".join(ex.map(_ocr_task, tasks))
    except Exception as e:
        print(f"[ERROR] PPTX 이미지 OCR 실패: {e}")
        return ""


def run_ocr_on_hwpx_images(hwpx_bytes: bytes) -> str:
    if reader is None:
        return ""
    try:
        with zipfile.ZipFile(io.BytesIO(hwpx_bytes)) as z:
            tasks = [(n, z.read(n)) for n in z.namelist() if n.startswith("Contents/") and n.lower().endswith((".jpg",".png",".bmp",".jpeg",".gif"))]
            if not tasks:
                return ""
            with ThreadPoolExecutor() as ex:
                return "\n".join(ex.map(_ocr_task, tasks))
    except Exception as e:
        print(f"[ERROR] HWPX 이미지 OCR 실패: {e}")
        return ""

# ==========================
# 파일 파싱
# ==========================

def parse_file(File_Bytes: bytes, File_Ext: str) -> tuple:
    File_Ext = (File_Ext or '').lower()

    if File_Ext == "txt":
        try:
            return File_Bytes.decode("utf-8"), False
        except UnicodeDecodeError:
            return File_Bytes.decode("cp949", errors='ignore'), False

    elif File_Ext == "doc":
        if win32com is None:
            raise ValueError("[ERROR] .doc 파싱을 위해서는 Windows 환경에 MS Office 설치가 필요합니다.")
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp:
                tmp.write(File_Bytes)
                tmp_path = tmp.name
            
            tmp_docx_path = tmp_path + "x"
            word = win32com.client.Dispatch("Word.Application")
            doc = word.Documents.Open(tmp_path)
            doc.SaveAs(tmp_docx_path, FileFormat=16)
            doc.Close()
            word.Quit()
            
            with open(tmp_docx_path, "rb") as f:
                converted_bytes = f.read()
                
            os.remove(tmp_path)
            os.remove(tmp_docx_path)
            
            return parse_file(converted_bytes, "docx")
        except Exception as e:
            raise ValueError(f"[ERROR] DOC -> DOCX 변환 실패: {e}")

    elif File_Ext == "docx":
        if Document is None:
            raise ValueError("[ERROR] python-docx 라이브러리 미설치")
        try:
            doc = Document(io.BytesIO(File_Bytes))
            text = "\n".join([p.text for p in doc.paragraphs if p.text])
            for table in doc.tables:
                for row in table.rows:
                    text += "\n" + " ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
            if not text.strip():
                text = run_ocr_on_docx_images(File_Bytes)
            return re.sub(r'\s+', ' ', text.strip()), False
        except Exception as e:
            raise ValueError(f"[ERROR] DOCX 파싱 실패: {e}")

    elif File_Ext == "pdf":
        if fitz is None:
            raise ValueError("[ERROR] PyMuPDF(fitz) 라이브러리 미설치")
        try:
            doc = fitz.open(stream=File_Bytes, filetype="pdf")
            if doc.is_encrypted:
                raise ValueError("[ERROR] 암호화된 PDF 문서")
            
            # 병렬 처리로 텍스트 추출
            def extract_page_text(page_num):
                return doc[page_num].get_text().replace("\n", " ")
            
            page_count = len(doc)
            if page_count > 1:
                max_workers = min(8, os.cpu_count() or 4)
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    texts = list(executor.map(extract_page_text, range(page_count)))
                text = " ".join(texts)
            else:
                text = doc[0].get_text().replace("\n", " ")
            
            if not any(ch.isalnum() for ch in text):
                text = run_ocr_on_pdf_images(File_Bytes)
            return text.strip(), False
        except Exception as e:
            raise ValueError(f"[ERROR] PDF 파싱 실패: {e}")

    elif File_Ext == "hwp":
        if File_Bytes[:4] == b'PK\x03\x04':
            print("[INFO] HWPX 파일로 감지됨, HWPX 파싱으로 전환")
            return parse_file(File_Bytes, "hwpx")
        
        if olefile is None:
            raise ValueError("[ERROR] olefile 라이브러리 미설치")
        try:
            ole = olefile.OleFileIO(io.BytesIO(File_Bytes))
            text = ""
            if ole.exists("PrvText"):
                text = ole.openstream("PrvText").read().decode("utf-16", errors="ignore").strip()
            if not text.strip():
                # BodyText, HeaderText, FooterText 모두 추출
                for entry in ole.listdir():
                    if entry[0] in ["BodyText", "HeaderText", "FooterText"]:
                        try:
                            raw = ole.openstream(entry).read()
                            text += ''.join(c for c in raw.decode("utf-16", errors="ignore") if c.isprintable() or c in '\n\r\t ') + "\n"
                        except Exception:
                            continue
            # 텍스트 유무와 관계없이 항상 OCR 실행 (이미지 내 텍스트 탐지)
            ocr_text = run_ocr_on_hwp_images(File_Bytes)
            if ocr_text:
                print(f"[INFO] HWP 이미지 OCR 추출: {len(ocr_text)}글자")
                text = (text + "\n" + ocr_text).strip()
            ole.close()
            text = re.sub(r'\s+', ' ', text).strip()
            return text, False
        except Exception as e:
            raise ValueError(f"[ERROR] HWP(olefile) 파싱 실패: {e}")

    elif File_Ext == "hwpx":
        try:
            with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                sections = [n for n in z.namelist() if n.startswith('Contents/section')]
                if not sections:
                    xmls = [n for n in z.namelist() if n.endswith('.xml')]
                    text = ''
                    for name in xmls:
                        data = z.read(name).decode('utf-8', errors='ignore')
                        text += re.sub('<[^>]+>', ' ', data)
                else:
                    text = ''
                    for name in sections:
                        data = z.read(name).decode('utf-8', errors='ignore')
                        text += re.sub('<[^>]+>', ' ', data)
                text = re.sub(r'\s+', ' ', text).strip()
                if not text:
                    text = run_ocr_on_hwpx_images(File_Bytes)
                return text.strip(), False
        except Exception as e:
            raise ValueError(f"[ERROR] HWPX 파싱 실패: {e}")
            
    elif File_Ext == "xlsx":
        if load_workbook is None:
            raise ValueError("[ERROR] openpyxl 라이브러리 미설치")
        try:
            # use read_only mode for robustness and memory
            wb = load_workbook(io.BytesIO(File_Bytes), data_only=True, read_only=True)
            lines = []
            for sheet in wb.worksheets:
                for ridx, row in enumerate(sheet.iter_rows(values_only=True)):
                    try:
                        # safely convert each cell to string; handle unexpected cell types
                        cells = []
                        for c in row:
                            if c is None:
                                continue
                            if isinstance(c, (list, tuple)):
                                # rare case: cell contains list-like value
                                cells.append(" ".join(map(str, c)))
                            else:
                                cells.append(str(c))
                        if cells:
                            lines.append(" ".join(cells))
                    except Exception as inner_e:
                        logging.warning(f"[WARN] XLSX row parse skipped: sheet={getattr(sheet,'title',None)} row={ridx} error={inner_e}")
                        continue
            text = "\n".join(lines).strip()
            return re.sub(r'\s+', ' ', text), False
        except Exception as e:
            # attempt a low-level zip/xml fallback to salvage text from sharedStrings/sheets
            try:
                with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                    s = []
                    # try sharedStrings
                    if 'xl/sharedStrings.xml' in z.namelist():
                        ss_xml = z.read('xl/sharedStrings.xml').decode('utf-8', errors='ignore')
                        parts = re.findall(r'<t[^>]*>(.*?)</t>', ss_xml, flags=re.DOTALL)
                        for p in parts:
                            s.append(re.sub(r'\s+', ' ', p.strip()))
                    # try sheet xmls
                    sheet_names = [n for n in z.namelist() if n.startswith('xl/worksheets/sheet') and n.endswith('.xml')]
                    for name in sheet_names:
                        try:
                            xml = z.read(name).decode('utf-8', errors='ignore')
                            vals = re.findall(r'<v>(.*?)</v>', xml, flags=re.DOTALL)
                            for v in vals:
                                s.append(re.sub(r'\s+', ' ', v.strip()))
                        except Exception:
                            continue
                    if not s:
                        raise Exception('xlsx low-level fallback produced no text')
                    return re.sub(r'\s+', ' ', '\n'.join(s).strip()), False
            except Exception as e2:
                raise ValueError(f"[ERROR] XLSX 파싱 실패: {e} | fallback: {e2}")

    elif File_Ext == "xls":
        if xlrd is None:
            if win32com is None:
                raise ValueError("[ERROR] .xls 파일을 처리하려면 xlrd 또는 win32com 라이브러리가 필요합니다.")
            
            print("[WARN] xlrd 라이브러리가 없어 win32com으로 .xls 파일을 처리합니다. (Windows/MS Office 환경 필요)")
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".xls") as tmp:
                    tmp.write(File_Bytes); tmp_path = tmp.name
                xlsx_path = tmp_path + "x"
                excel = win32com.client.Dispatch("Excel.Application")
                wb = excel.Workbooks.Open(tmp_path)
                wb.SaveAs(xlsx_path, FileFormat=51); wb.Close(SaveChanges=False); excel.Quit()
                with open(xlsx_path, "rb") as f:
                    converted_bytes = f.read()
                os.remove(tmp_path); os.remove(xlsx_path)
                return parse_file(converted_bytes, "xlsx")
            except Exception as e:
                raise ValueError(f"[ERROR] win32com을 이용한 XLS → XLSX 변환 실패: {e}")
        try:
            workbook = xlrd.open_workbook(file_contents=File_Bytes)
            text = ""
            for sheet in workbook.sheets():
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
            return text.strip(), False
        except Exception as e:
            raise ValueError(f"[ERROR] XLS (xlrd) 파싱 실패: {e}")

    elif File_Ext in ["ppt","pptx"]:
        if File_Ext == "ppt":
            if win32com is None:
                raise ValueError("[ERROR] PPT 파싱은 Windows/MS Office 환경 필요")
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
                    tmp.write(File_Bytes); tmp_path = tmp.name
                pptx_path = tmp_path + "x"
                pp = win32com.client.Dispatch("PowerPoint.Application")
                pres = pp.Presentations.Open(tmp_path, WithWindow=False)
                pres.SaveAs(pptx_path, FileFormat=24); pres.Close(); pp.Quit()
                with open(pptx_path, "rb") as f:
                    File_Bytes = f.read()
                os.remove(tmp_path); os.remove(pptx_path)
            except Exception as e:
                raise ValueError(f"[ERROR] PPT → PPTX 변환 실패: {e}")
        if Presentation is None:
            raise ValueError("[ERROR] python-pptx 라이브러리 미설치")
        try:
            prs = Presentation(io.BytesIO(File_Bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
            ocr = run_ocr_on_pptx_images(File_Bytes)
            if ocr:
                text += "\n" + ocr
            return text.strip(), False
        except Exception as e:
            raise ValueError(f"[ERROR] PPTX 파싱 실패: {e}")
            
    # [수정] GIF 파일 처리 로직 추가
    elif File_Ext == "gif":
        if ImageSequence is None:
            raise ValueError("[ERROR] GIF 처리를 위해 Pillow(PIL) 라이브러리가 필요합니다.")
        print("[INFO] GIF 파일 감지: 다중 프레임 OCR 시작")
        try:
            img = Image.open(io.BytesIO(File_Bytes))
            ocr_text = ""
            frame_count = 0
            for i, frame in enumerate(ImageSequence.Iterator(img)):
                frame_count += 1
                if i % 3 != 0:  # 3프레임 마다 샘플링
                    continue
                frame_rgb = frame.convert("RGB")
                result = reader.readtext(np.array(frame_rgb))
                for box in result:
                    ocr_text += box[1] + "\n"
                if len(ocr_text) > 1000: # 텍스트가 일정 길이 이상이면 조기 종료
                    break
            print(f"[INFO] GIF OCR 완료: {frame_count}프레임 중 {len(ocr_text)}글자 추출")
            return ocr_text.strip(), True
        except Exception as e:
            print(f"[ERROR] GIF OCR 실패: {e}")
            return "", True

    elif File_Ext in ["png","jpg","jpeg","bmp","webp","tiff"]:
        return run_ocr_on_single_image(File_Bytes), True

    else:
        raise ValueError(f"[ERROR] 지원하지 않는 파일 형식: {File_Ext}")

# ==========================
# 얼굴 탐지
# ==========================

def detect_faces_in_image_bytes(image_bytes, confidence_threshold=0.98):
    if detector is None or Image is None:
        return []
    try:
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        
        # 이미지 크기 체크 (너무 작으면 스킵)
        if img.width < 50 or img.height < 50:
            return []
        
        # 이미지 리사이즈 (큰 이미지는 축소)
        max_size = 800
        if img.width > max_size or img.height > max_size:
            r = min(max_size / img.width, max_size / img.height)
            img = img.resize((int(img.width*r), int(img.height*r)), Image.LANCZOS)
        
        img_np = np.array(img)
        results = detector.detect_faces(img_np)
        detections = []
        
        for res in results:
            conf = float(res.get('confidence', 0))
            x, y, w, h = res['box']
            
            # 필터링 조건
            # 1. confidence >= 0.98
            if conf < confidence_threshold:
                continue
            
            # 2. 얼굴 크기 검증 (너무 작거나 큰 것 제외)
            if w < 30 or h < 30 or w > img.width * 0.9 or h > img.height * 0.9:
                continue
            
            # 3. 가로세로 비율 검증 (0.6~1.5)
            aspect_ratio = w / h if h > 0 else 0
            if aspect_ratio < 0.6 or aspect_ratio > 1.5:
                continue
            
            # 4. keypoints 검증 (눈, 코, 입)
            keypoints = res.get('keypoints', {})
            if keypoints:
                left_eye = keypoints.get('left_eye')
                right_eye = keypoints.get('right_eye')
                nose = keypoints.get('nose')
                
                # 눈과 코가 모두 탐지되어야 함
                if not (left_eye and right_eye and nose):
                    continue
                
                # 두 눈 사이 거리 검증
                eye_distance = abs(left_eye[0] - right_eye[0])
                if eye_distance < w * 0.2 or eye_distance > w * 0.8:
                    continue
            
            detections.append({"bbox":[int(x),int(y),int(w),int(h)], "confidence": conf})
        
        return detections
    except Exception:
        return []


def _face_task(args):
    name, bytes_ = args
    faces = detect_faces_in_image_bytes(bytes_)
    return {"image_name": name, "faces_found": len(faces), "faces": faces} if faces else None


def scan_file_for_face_images(file_bytes, file_ext):
    file_ext = (file_ext or '').lower()
    tasks = []
    try:
        if file_ext in ["png","jpg","jpeg","bmp","webp","gif","tiff"]:
            tasks.append(("uploaded_image", file_bytes))
        elif file_ext in ["docx","pptx","hwpx"]:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                if file_ext == 'docx':
                    prefix = 'word/media/'
                    imgs = [n for n in z.namelist() if n.startswith(prefix)]
                elif file_ext == 'pptx':
                    prefix = 'ppt/media/'
                    imgs = [n for n in z.namelist() if n.startswith(prefix)]
                else:  # hwpx
                    imgs = [n for n in z.namelist() if n.startswith('Contents/') and n.lower().endswith((".png",".jpg",".jpeg",".bmp",".gif"))]
                tasks.extend((n, z.read(n)) for n in imgs)
        elif file_ext == 'pdf' and fitz:
            doc = fitz.open(stream=file_bytes, filetype='pdf')
            seen = set()
            for p, page in enumerate(doc):
                for i, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    if xref in seen:
                        continue
                    seen.add(xref)
                    bi = doc.extract_image(xref)
                    if not bi or 'image' not in bi or len(bi['image']) < 5000:
                        continue
                    tasks.append((f"pdf_p{p+1}_img{i+1}", bi['image']))
        elif file_ext == 'hwp' and olefile:
            ole = olefile.OleFileIO(io.BytesIO(file_bytes))
            tasks.extend(("/".join(e), ole.openstream(e).read()) for e in ole.listdir() if e[0]=='BinData')
    except Exception as e:
        print(f"[WARN] {file_ext.upper()} 이미지 추출 실패: {e}")
    if not tasks:
        return []
    with ThreadPoolExecutor(max_workers=min(8, (os.cpu_count() or 4))) as ex:
        return [r for r in ex.map(_face_task, tasks) if r]

# ==========================
# 메인 핸들러
# ==========================

def handle_input_raw(Input_Data: bytes, Original_Format: str = None, Original_Filename: str = None):
    if not isinstance(Input_Data, bytes):
        raise ValueError("지원하지 않는 입력 형식입니다.")
    print(f"\n[INFO] ========== 파일 처리 시작 (확장자: {Original_Format}) ==========")

    # 병렬 처리: 텍스트 추출과 얼굴 탐지를 동시에 실행
    Parsed_Text = ""
    is_image_only = False
    parse_error = None
    with ThreadPoolExecutor(max_workers=2) as executor:
        text_future = executor.submit(parse_file, Input_Data, Original_Format or "")
        face_future = executor.submit(scan_file_for_face_images, Input_Data, Original_Format or "")

        # 얼굴 탐지 결과는 가능하면 항상 확보
        try:
            image_detections = face_future.result()
        except Exception as e:
            logging.warning(f"파일 이미지(얼굴) 추출 중 오류: {e}")
            image_detections = []

        # 텍스트 파싱에서 오류가 발생하면 예외를 잡아 전달 가능한 형태로 기록
        try:
            Parsed_Text, is_image_only = text_future.result()
        except Exception as e:
            parse_error = e
            Parsed_Text = ""
            is_image_only = False
            logging.warning(f"파일 파싱 실패: {e}")
    
    print(f"[INFO] 추출된 텍스트 길이: {len(Parsed_Text)} 글자")
    if Parsed_Text:
        print(f"[INFO] 텍스트 미리보기: {Parsed_Text[:200]}...")
    else:
        print(f"[WARN] 추출된 텍스트 없음!")

    Detected = []
    comb = None

    combined_text = Parsed_Text or ""
    if Original_Filename:
        base = Original_Filename.rsplit('.', 1)[0]
        combined_text = (base + " \n" + combined_text).strip()

    if combined_text:
        ner_results   = detect_by_ner(combined_text)
        regex_results = detect_by_regex(combined_text)
        quasi_results = detect_quasi_identifiers(combined_text)
        
        all_detected = regex_results + ner_results + quasi_results
        
        face_items_for_risk = [{"type": "image_face", "value": "얼굴사진"}] * len(image_detections)
        final_all_detected = all_detected + face_items_for_risk
        
        comb = analyze_combination_risk(final_all_detected, combined_text)

        # NOTE: 조합위험(combination_risk)은 내부 메타데이터로 유지하되
        # 탐지 결과 목록(Detected)에는 추가하지 않습니다. 호출자에서
        # 별도의 필드로 받아 처리하도록 반환값으로 comb를 포함합니다.
        Detected = all_detected

    total_faces = 0
    for img in image_detections:
        cnt = img.get('faces_found', 0)
        total_faces += cnt
        if cnt > 0:
            Detected.append({"type":"image_face","value":f"{img.get('image_name','이미지')} 내 얼굴 {cnt}개","detail":img})
    if total_faces > 0:
        print(f"[INFO] [OK] 이미지 얼굴 총 {total_faces}개 탐지")
    


    masked_filename = None
    if Original_Filename:
        masked, types = mask_pii_in_filename(Original_Filename)
        masked_filename = masked if masked != Original_Filename else None

    # 파싱 에러가 발생했을 경우, 탐지 항목에 오류로 남기고 backend_status를 False로 설정
    if parse_error is not None:
        err_msg = str(parse_error)
        Detected.append({
            "type": "file_parse_error",
            "value": err_msg,
            "detail": {"filename": Original_Filename, "format": Original_Format}
        })
        backend_status = False
    else:
        backend_status = True

    return Detected, (masked_filename or ""), backend_status, image_detections, comb