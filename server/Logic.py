import io
import re
import logging
import numpy as np
import datetime
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None
try:
    import easyocr
except ImportError:
    easyocr = None
import zipfile
try:
    from PIL import Image
except ImportError:
    Image = None
from transformers import AutoTokenizer, AutoModelForTokenClassification, pipeline
import os
try:
    from docx import Document
except ImportError:
    Document = None
try:
    from mtcnn import MTCNN
except ImportError:
    MTCNN = None
try:
    import olefile
except ImportError:
    olefile = None
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    import win32com.client
except ImportError:
    win32com = None

logging.basicConfig(level=logging.INFO, format='[%(levelname)s] %(message)s')

# 허깅페이스 토큰 설정 (환경 변수에서 읽기)
HF_TOKEN = os.getenv("HF_TOKEN", None)

NER_MODEL_NAME = "soddokayo/klue-roberta-base-ner"
print(f"[INFO] NER 모델 로딩 중: {NER_MODEL_NAME}")

if HF_TOKEN:
    ner_tokenizer = AutoTokenizer.from_pretrained(NER_MODEL_NAME, token=HF_TOKEN)
    ner_model = AutoModelForTokenClassification.from_pretrained(NER_MODEL_NAME, token=HF_TOKEN)
    print("[INFO] ✓ 허깅페이스 토큰 인증 완료")
else:
    print("[WARN] HF_TOKEN 환경 변수 없음 - 공개 모델로 시도")
    ner_tokenizer = AutoTokenizer.from_pretrained(NER_MODEL_NAME)
    ner_model = AutoModelForTokenClassification.from_pretrained(NER_MODEL_NAME)

ner_pipeline = pipeline("ner", model=ner_model, tokenizer=ner_tokenizer, grouped_entities=True)
print("[INFO] ✓ NER 모델 로딩 완료")

try:
    reader = easyocr.Reader(['ko', 'en'], gpu=False)
    print("[INFO] EasyOCR 초기화 완료")
except Exception as e:
    print(f"[WARN] EasyOCR 초기화 실패: {e}")
    reader = None

detector = MTCNN()

# ==========================
# 검증 함수 (Luhn + 주민번호)
# ==========================
def validate_luhn(card_number: str) -> bool:
    """Luhn 알고리즘으로 카드번호 검증"""
    digits = [int(d) for d in re.sub(r"\D", "", card_number)]
    checksum = 0
    for i, d in enumerate(reversed(digits)):
        if i % 2 == 1:
            t = d * 2
            checksum += t - 9 if t > 9 else t
        else:
            checksum += d
    return checksum % 10 == 0

def validate_ssn(ssn: str) -> bool:
    """주민등록번호 검증 (생년월일 + 체크섬)"""
    ssn = re.sub(r"\D", "", ssn)
    if len(ssn) != 13:
        return False
    yy = int(ssn[0:2])
    mm = int(ssn[2:4])
    dd = int(ssn[4:6])
    g = int(ssn[6])
    
    if g not in [1, 2, 3, 4, 5, 6, 7, 8]:
        return False
    
    full_year = (1900 + yy) if g in [1, 2, 5, 6] else (2000 + yy)
    
    try:
        datetime.date(full_year, mm, dd)
    except ValueError:
        return False
    
    if full_year >= 2020:
        return True
    
    weights = [2, 3, 4, 5, 6, 7, 8, 9, 2, 3, 4, 5]
    s = sum(int(ssn[i]) * weights[i] for i in range(12))
    check = (11 - (s % 11)) % 10
    return check == int(ssn[-1])

def validate_driver_license(license_num: str) -> bool:
    """운전면허번호 검증 (지역코드 체크)"""
    license_num = re.sub(r"\D", "", license_num)
    if len(license_num) != 12:
        return False
    
    # 유효한 지역코드 (11~26, 24 제외)
    valid_region_codes = {
        '11', '12', '13', '14', '15', '16', '17', '18', '19',
        '20', '21', '22', '23', '25', '26'
    }
    
    region_code = license_num[:2]
    return region_code in valid_region_codes

# ==========================
# 파일 파싱
# ==========================
def parse_file(File_Bytes: bytes, File_Ext: str) -> tuple:
    """
    파일을 파싱하여 텍스트를 추출합니다.
    Returns: (text, is_pure_image)
    """
    File_Ext = File_Ext.lower()
    
    if File_Ext == "txt":
        try:
            return File_Bytes.decode("utf-8"), False
        except UnicodeDecodeError:
            return File_Bytes.decode("cp949"), False

    elif File_Ext == "docx":
        try:
            file_stream = io.BytesIO(File_Bytes)
            doc = Document(file_stream)
        except Exception:
            raise ValueError("[ERROR] DOCX 파싱 실패: Document 객체 생성 불가")
        try:
            # 단락 텍스트 추출
            print(f"[DEBUG] DOCX 단락 개수: {len(doc.paragraphs)}")
            para_texts = []
            for i, para in enumerate(doc.paragraphs):
                if para.text:
                    print(f"[DEBUG] 단락 {i}: '{para.text[:50]}'")
                    para_texts.append(para.text)
            text = "\n".join(para_texts)
            
            # 테이블 텍스트 추출
            print(f"[DEBUG] DOCX 테이블 개수: {len(doc.tables)}")
            for t_idx, table in enumerate(doc.tables):
                print(f"[DEBUG] 테이블 {t_idx}: {len(table.rows)}행 x {len(table.columns)}열")
                for r_idx, row in enumerate(table.rows):
                    row_text = " ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        print(f"[DEBUG] 테이블 {t_idx} 행 {r_idx}: '{row_text[:50]}'")
                        text += "\n" + row_text
            
            print(f"[DEBUG] 추출된 텍스트 길이: {len(text)} 글자")
            if not text.strip():
                print("[INFO] 텍스트 없음 → OCR 실행")
                text = run_ocr_on_docx_images(File_Bytes)
                if not text.strip():
                    print("[WARN] OCR 텍스트 추출 실패, 빈 문자열 반환")
            return text, False
        except Exception as e:
            raise ValueError(f"[ERROR] DOCX 파싱 실패 (내용 추출 중 오류): {e}")

    elif File_Ext == "pdf":
        try:
            doc = fitz.open(stream=File_Bytes, filetype="pdf")
            if doc.is_encrypted:
                raise ValueError("[ERROR] PDF 파싱 실패: 암호로 보호된 PDF 문서입니다.")
            
            # 텍스트 추출 최적화 (리스트 컴프리헨션 대신 join 사용)
            text_parts = []
            for page in doc:
                page_text = page.get_text()
                if page_text:
                    text_parts.append(page_text.replace("\n", " "))
            text = " ".join(text_parts)
            
            if not any(char.isalnum() for char in text):
                print("[INFO] 텍스트 없음 → OCR 실행")
                text = run_ocr_on_pdf_images(File_Bytes)
                if not text.strip():
                    print("[WARN] OCR 텍스트 추출 실패, 빈 문자열 반환")
            return text, False
        except Exception as e:
            raise ValueError(f"[ERROR] PDF 파싱 실패: {e}")
    
    elif File_Ext == "hwp":
        # HWPX 파일 자동 감지 (ZIP 시그니처 체크)
        if File_Bytes[:4] == b'PK\x03\x04':  # ZIP 파일 시그니처
            print("[INFO] HWPX 파일로 감지됨, HWPX 파싱으로 전환")
            return parse_file(File_Bytes, "hwpx")
        
        # win32com으로 HWP → HWPX 변환 시도
        if win32com is not None:
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".hwp") as tmp:
                    tmp.write(File_Bytes)
                    tmp_path = tmp.name
                tmp_hwpx_path = tmp_path + "x"
                try:
                    hwp = win32com.client.Dispatch("HWPFrame.HwpObject")
                    hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModule")
                    hwp.Open(tmp_path)
                    hwp.SaveAs(tmp_hwpx_path, "HWPX")
                    hwp.Quit()
                except Exception:
                    raise ValueError("[ERROR] HWP 변환 실패 또는 암호화")
                with open(tmp_hwpx_path, "rb") as f:
                    converted = f.read()
                os.remove(tmp_path)
                os.remove(tmp_hwpx_path)
                return parse_file(converted, "hwpx")
            except Exception as e:
                print(f"[WARN] HWP → HWPX 변환 실패, 기본 로직으로 진행: {e}")
        
        # 기본 HWP 파싱
        if olefile is None:
            raise ValueError("[ERROR] HWP 파싱 실패: olefile 라이브러리 미설치")
        try:
            ole = olefile.OleFileIO(io.BytesIO(File_Bytes))
            text = ""
            
            # 방법 1: PrvText 스트림에서 추출
            if ole.exists("PrvText"):
                stream = ole.openstream("PrvText")
                raw = stream.read()
                text = raw.decode("utf-16", errors="ignore").strip()
                print(f"[DEBUG] HWP PrvText 추출: {len(text)} 글자")
            
            # 방법 2: BodyText 섹션에서 추출 (더 정확함)
            if not text.strip():
                print("[INFO] PrvText 없음, BodyText 시도")
                for entry in ole.listdir():
                    if entry[0] == "BodyText":
                        try:
                            stream = ole.openstream(entry)
                            raw = stream.read()
                            decoded = raw.decode("utf-16", errors="ignore")
                            cleaned = ''.join(c for c in decoded if c.isprintable() or c in '\n\r\t ')
                            text += cleaned + "\n"
                        except Exception:
                            continue
                print(f"[DEBUG] HWP BodyText 추출: {len(text)} 글자")
            
            # 방법 3: OCR 시도
            if not text.strip():
                print("[INFO] HWP 텍스트 없음 → 이미지 OCR 시도")
                text = run_ocr_on_hwp_images(File_Bytes)
            
            ole.close()
            text = text.strip()
            
            # HWP 텍스트 정리 (태그만 제거, 내용은 보존)
            text = re.sub(r'<>', ' ', text)  # 빈 태그 제거
            text = re.sub(r'<', ' ', text)  # < 기호를 공백으로
            text = re.sub(r'>', ' ', text)  # > 기호를 공백으로
            text = re.sub(r'\s+', ' ', text)  # 여러 공백을 하나로
            text = text.strip()
            
            # 한글 띄어쓰기 정규화 ("이 무 송" -> "이무송")
            original_text = text
            while True:
                prev = text
                text = re.sub(r'([가-힣])\s+([가-힣])(?=\s|[가-힣]|$)', r'\1\2', text)
                if prev == text:
                    break
            
            if text != original_text:
                print(f"[INFO] ⚠ HWP 한글 띄어쓰기 정규화 적용")
                print(f"[DEBUG] 원본: {original_text[:200]}")
                print(f"[DEBUG] 정규화: {text[:200]}")
            
            print(f"[INFO] HWP 최종 텍스트: {len(text)} 글자")
            print(f"[DEBUG] HWP 텍스트 미리보기: {text[:200]}")
            return text, False
        except Exception as e:
            error_msg = str(e)
            if "not an OLE2" in error_msg:
                raise ValueError("[ERROR] HWP 파싱 실패: 암호화된 파일이거나 손상된 파일입니다. 암호를 해제하거나 다른 파일을 사용해주세요.")
            raise ValueError(f"[ERROR] HWP 파싱 실패: {e}")
    
    elif File_Ext == "hwpx":
        try:
            with zipfile.ZipFile(io.BytesIO(File_Bytes)) as z:
                xml_files = [n for n in z.namelist() if n.endswith('.xml')]
                if not xml_files:
                    raise ValueError("[ERROR] HWPX 파싱 실패: 암호 또는 손상")
                text = ""
                for name in xml_files:
                    data = z.read(name).decode("utf-8", errors="ignore")
                    text += re.sub("<[^>]+>", " ", data)
                text = text.strip()
                if not text:
                    print("[INFO] HWPX 텍스트 없음 → OCR 실행")
                    text = run_ocr_on_hwpx_images(File_Bytes)
                return text, False
        except Exception as e:
            raise ValueError(f"[ERROR] HWPX 파싱 실패: {e}")
    
    elif File_Ext in ["xls", "xlsx"]:
        if File_Ext == "xlsx":
            if load_workbook is None:
                raise ValueError("[ERROR] XLSX 파싱 실패: openpyxl 라이브러리 미설치")
            try:
                wb = load_workbook(io.BytesIO(File_Bytes), data_only=True)
                text = ""
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join([str(cell).strip() for cell in row if cell is not None and str(cell).strip()])
                        if row_text:
                            text += row_text + ". "  # 행 구분자 추가
                return text.strip(), False
            except Exception as e:
                raise ValueError(f"[ERROR] XLSX 파싱 실패: {e}")
        else:  # .xls
            if win32com is None:
                raise ValueError("[ERROR] XLS 파싱 실패: win32com 미설치")
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".xls") as tmp:
                    tmp.write(File_Bytes)
                    tmp_path = tmp.name
                tmp_xlsx_path = tmp_path + "x"
                try:
                    excel = win32com.client.Dispatch("Excel.Application")
                    excel.Visible = False
                    wb = excel.Workbooks.Open(tmp_path)
                    wb.SaveAs(tmp_xlsx_path, FileFormat=51)
                    wb.Close(SaveChanges=False)
                    excel.Quit()
                except Exception:
                    raise ValueError("[ERROR] XLS → XLSX 변환 실패 또는 암호")
                with open(tmp_xlsx_path, "rb") as f:
                    converted = f.read()
                os.remove(tmp_path)
                os.remove(tmp_xlsx_path)
                return parse_file(converted, "xlsx")
            except Exception as e:
                raise ValueError(f"[ERROR] XLS 파싱 실패: {e}")
    
    elif File_Ext == "pptx":
        if Presentation is None:
            raise ValueError("[ERROR] PPTX 파싱 실패: python-pptx 라이브러리 미설치")
        try:
            prs = Presentation(io.BytesIO(File_Bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            ocr_text = run_ocr_on_pptx_images(File_Bytes)
            if ocr_text:
                text += "\n" + ocr_text
            return text.strip(), False
        except Exception as e:
            raise ValueError(f"[ERROR] PPTX 파싱 실패: {e}")
    
    elif File_Ext == "ppt":
        raise ValueError("[ERROR] PPT 파일은 지원하지 않습니다. PPTX로 변환 후 업로드해주세요.")
    
    elif File_Ext == "doc":
        if win32com is None:
            raise ValueError("[ERROR] DOC 파싱 실패: win32com 미설치")
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp:
                tmp.write(File_Bytes)
                tmp_path = tmp.name
            tmp_docx_path = tmp_path + "x"
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(tmp_path)
                doc.SaveAs(tmp_docx_path, FileFormat=16)
                doc.Close()
                word.Quit()
            except Exception:
                raise ValueError("[ERROR] DOC 변환/파싱 실패: 암호")
            with open(tmp_docx_path, "rb") as f:
                converted = f.read()
            os.remove(tmp_path)
            os.remove(tmp_docx_path)
            return parse_file(converted, "docx")
        except Exception as e:
            raise ValueError(f"[ERROR] DOC 파싱 실패: {e}")
    
    elif File_Ext in ["png", "jpg", "jpeg", "bmp", "webp", "gif", "tiff"]:
        print(f"[INFO] 이미지 파일 감지: {File_Ext}")
        ocr_text = run_ocr_on_single_image(File_Bytes)
        return ocr_text, True
    
    else:
        raise ValueError(f"[ERROR] 지원하지 않는 파일 형식: {File_Ext}")

# ==========================
# OCR
# ==========================
def run_ocr_on_single_image(image_bytes: bytes) -> str:
    """단일 이미지에서 OCR 수행"""
    if reader is None:
        print("[WARN] EasyOCR이 초기화되지 않았습니다.")
        return ""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        result = reader.readtext(np.array(img))
        ocr_text = "\n".join([box[1] for box in result])
        print(f"[INFO] 이미지 OCR 완료: {len(ocr_text)} 글자 추출")
        return ocr_text.strip()
    except Exception as e:
        print(f"[ERROR] 이미지 OCR 실패: {e}")
        return ""

def run_ocr_on_docx_images(file_bytes):
    if reader is None:
        print("[WARN] EasyOCR이 초기화되지 않았습니다.")
        return ""
    ocr_text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as docx_zip:
            image_files = [f for f in docx_zip.namelist() if f.startswith("word/media/")]
            print(f"[INFO] DOCX 내부 이미지 {len(image_files)}개 발견, OCR 시작")
            for image_name in image_files:
                try:
                    with docx_zip.open(image_name) as image_file:
                        image = Image.open(image_file).convert('RGB')
                        print(f"[DEBUG] 이미지 크기: {image.size}")
                        
                        # 이미지 전처리: 대비 증가
                        from PIL import ImageEnhance
                        enhancer = ImageEnhance.Contrast(image)
                        image = enhancer.enhance(2.0)
                        
                        # OCR 실행 (신뢰도 임계값 낮춤)
                        result = reader.readtext(np.array(image), detail=1, paragraph=False)
                        print(f"[DEBUG] OCR 결과: {len(result)}개 텍스트 박스")
                        for box in result:
                            text_content = box[1]
                            confidence = box[2]
                            print(f"[DEBUG] 추출 텍스트: '{text_content}' (신뢰도: {confidence:.2f})")
                            if confidence > 0.1:  # 신뢰도 10% 이상만
                                ocr_text += text_content + "\n"
                except Exception as img_e:
                    print(f"[ERROR] 이미지 {image_name} OCR 실패: {img_e}")
    except Exception as e:
        print(f"[ERROR] DOCX 이미지 OCR 실패: {e}")
    return ocr_text.strip()

def run_ocr_on_pdf_images(pdf_bytes: bytes) -> str:
    ocr_text = ""
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        print(f"[INFO] PDF {len(doc)}페이지 OCR 시작")
        for idx, page in enumerate(doc):
            pix = page.get_pixmap()
            img = Image.open(io.BytesIO(pix.tobytes("ppm")))
            try:
                result = reader.readtext(np.array(img))
                for box in result:
                    ocr_text += box[1] + "\n"
            except Exception as e:
                print(f"[ERROR] EasyOCR 실패 (페이지 {idx}): {e}")
    return ocr_text.strip()

def run_ocr_on_hwp_images(hwp_bytes: bytes) -> str:
    """HWP 파일 내부 이미지에서 OCR 수행"""
    if reader is None or olefile is None:
        return ""
    ocr_text = ""
    try:
        ole = olefile.OleFileIO(io.BytesIO(hwp_bytes))
        # HWP 이미지는 BinData 스트림에 저장됨
        for entry in ole.listdir():
            if entry[0] == "BinData":
                try:
                    stream = ole.openstream(entry)
                    img_bytes = stream.read()
                    img = Image.open(io.BytesIO(img_bytes))
                    result = reader.readtext(np.array(img))
                    for box in result:
                        ocr_text += box[1] + "\n"
                except Exception:
                    continue
        ole.close()
        if ocr_text:
            print(f"[INFO] HWP 이미지 OCR 완료: {len(ocr_text)} 글자 추출")
    except Exception as e:
        print(f"[ERROR] HWP 이미지 OCR 실패: {e}")
    return ocr_text.strip()

def run_ocr_on_pptx_images(pptx_bytes: bytes) -> str:
    """PPTX 파일 내부 이미지에서 OCR 수행"""
    if reader is None:
        return ""
    ocr_text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as pptx_zip:
            image_files = [f for f in pptx_zip.namelist() if f.startswith("ppt/media/")]
            print(f"[INFO] PPTX 내부 이미지 {len(image_files)}개 발견, OCR 시작")
            for image_name in image_files:
                try:
                    with pptx_zip.open(image_name) as image_file:
                        img = Image.open(image_file)
                        result = reader.readtext(np.array(img))
                        for box in result:
                            ocr_text += box[1] + "\n"
                except Exception:
                    continue
        if ocr_text:
            print(f"[INFO] PPTX 이미지 OCR 완료: {len(ocr_text)} 글자 추출")
    except Exception as e:
        print(f"[ERROR] PPTX 이미지 OCR 실패: {e}")
    return ocr_text.strip()

def run_ocr_on_hwpx_images(hwpx_bytes: bytes) -> str:
    """HWPX 파일 내부 이미지에서 OCR 수행"""
    if reader is None:
        return ""
    ocr_text = ""
    try:
        with zipfile.ZipFile(io.BytesIO(hwpx_bytes)) as z:
            for name in z.namelist():
                if name.startswith("Contents/") and name.lower().endswith((".jpg", ".png", ".bmp")):
                    img = Image.open(io.BytesIO(z.read(name)))
                    result = reader.readtext(np.array(img))
                    for box in result:
                        ocr_text += box[1] + "\n"
    except Exception as e:
        print(f"[ERROR] HWPX 이미지 OCR 실패: {e}")
    return ocr_text.strip()

# ==========================
# 정규식 탐지 (강화 - 띄어쓰기/하이픈 우회 방지)
# ==========================
def detect_by_regex(Text: str) -> list:
    """
    정규식 기반 개인정보 탐지 (우회 방지 강화)
    띄어쓰기, 하이픈 없이도 탐지 가능
    """
    # 탐지 전 텍스트 정규화 (공백/하이픈 제거한 버전도 함께 검사)
    normalized_text = re.sub(r'[\s\-]', '', Text)
    
    Patterns = {
        # 전화번호: 휴대폰(010-1234-5678) + 지역번호(02-1234-5678, 032 - 123 - 4567)
        "phone": re.compile(r"\b(01[016789]|0[2-6][0-9]?)[\s\-]*\d{3,4}[\s\-]*\d{4}\b"),
        
        # 이메일
        "email": re.compile(r"\b[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}\b"),
        
        # 생년월일 (1900~2006년생만)
        "birth": re.compile(r"\b(19[0-9]{2}|200[0-6])[년./\- ]+(0?[1-9]|1[0-2])[월./\- ]+(0?[1-9]|[12][0-9]|3[01])[일]?\b"),
        
        # 주민등록번호: 123456-1234567, 1234561234567 모두 탐지
        "ssn": re.compile(r"\b\d{6}[\s\-]?[1-4]\d{6}\b"),
        
        # 외국인등록번호
        "alien_reg": re.compile(r"\b\d{6}[\s\-]?[5-8]\d{6}\b"),
        
        # 운전면허번호
        "driver_license": re.compile(r"\b\d{2}[\s\-]?\d{2}[\s\-]?\d{6}[\s\-]?\d{2}\b"),
        
        # 여권번호
        "passport": re.compile(r"\b[A-Z]\d{2,3}[A-Z]?\d{4,5}\b"),
        
        # 계좌번호
        "account": re.compile(r"\b\d{6}[\s\-]?\d{2}[\s\-]?\d{6}\b"),
        
        # 카드번호: 1234-5678-9012-3456, 1234567890123456 모두 탐지
        "card": re.compile(r"\b(?:\d{4}[\s\-]?){3}\d{4}\b"),
        
        # IP 주소
        "ip": re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b")
    }
    
    detected = []
    
    # 원본 텍스트에서 탐지
    for label, pattern in Patterns.items():
        for match in pattern.finditer(Text):
            item = {"type": label, "value": match.group(), "span": match.span()}
            # 카드번호 Luhn 검증
            if label == "card":
                item["status"] = "valid" if validate_luhn(item["value"]) else "invalid (Luhn check failed)"
            # 주민번호 검증
            if label == "ssn":
                item["status"] = "valid" if validate_ssn(item["value"]) else "invalid (SSN check failed)"
            # 운전면허 검증
            if label == "driver_license":
                if not validate_driver_license(item["value"]):
                    continue  # 유효하지 않으면 스킵
            detected.append(item)
    
    # 정규화된 텍스트에서 추가 탐지 (공백/하이픈 우회 시도 차단)
    normalized_patterns = {
        # 전화번호 (공백/하이픈 없이): 휴대폰 + 지역번호
        "phone_normalized": re.compile(r"\b(01[016789]\d{7,8}|0[2-6][0-9]?\d{7,8})\b"),
        
        # 주민등록번호 (공백/하이픈 없이)
        "ssn_normalized": re.compile(r"\b\d{6}[1-4]\d{6}\b"),
        
        # 외국인등록번호 (공백/하이픈 없이)
        "alien_reg_normalized": re.compile(r"\b\d{6}[5-8]\d{6}\b"),
        
        # 운전면허번호 (공백/하이픈 없이)
        "driver_license_normalized": re.compile(r"\b\d{2}\d{2}\d{6}\d{2}\b"),
        
        # 계좌번호 (공백/하이픈 없이)
        "account_normalized": re.compile(r"\b\d{6}\d{2}\d{6}\b"),
        
        # 카드번호 (공백/하이픈 없이)
        "card_normalized": re.compile(r"\b\d{16}\b")
    }
    
    # 정규화된 텍스트에서 탐지 (중복 제거)
    existing_values = {d["value"].replace(" ", "").replace("-", "") for d in detected}
    
    for label, pattern in normalized_patterns.items():
        original_label = label.replace("_normalized", "")
        for match in pattern.finditer(normalized_text):
            normalized_value = match.group()
            if normalized_value not in existing_values:
                item = {"type": original_label, "value": normalized_value, "span": match.span()}
                # 카드번호 Luhn 검증
                if original_label == "card":
                    item["status"] = "valid" if validate_luhn(item["value"]) else "invalid (Luhn check failed)"
                # 주민번호 검증
                if original_label == "ssn":
                    item["status"] = "valid" if validate_ssn(item["value"]) else "invalid (SSN check failed)"
                # 운전면허 검증
                if original_label == "driver_license":
                    if not validate_driver_license(item["value"]):
                        continue  # 유효하지 않으면 스킵
                detected.append(item)
                existing_values.add(normalized_value)
    
    return detected

# ==========================
# NER 탐지 (조직명 제외)
# ==========================
# NER이 놓치는 흔한 이름 화이트리스트
NAME_WHITELIST = {'홍길동'}

def detect_by_ner(Text: str) -> list:
    if not Text.strip():
        return []
    
    # 정규식으로 이름 추출 ("성명이무송" 패턴)
    name_pattern = re.compile(r'성명([가-힣]{2,4})(?:생년월일|주소|연락처|전화|이메일|E-Mail|휴대폰|생년)')
    name_match = name_pattern.search(Text)
    extracted_name = None
    if name_match:
        extracted_name = name_match.group(1)
        print(f"[INFO] ✓ 정규식으로 이름 추출: {extracted_name}")
    
    print(f"[DEBUG] NER 입력: {Text[:100]}")
    Results = ner_pipeline(Text)
    print(f"[DEBUG] NER 원본 결과: {Results}")
    
    Detected = []
    location_parts = []  # LC(주소) 병합용
    
    # 정규식으로 추출한 이름 먼저 추가
    if extracted_name:
        Detected.append({"type": "PS", "value": extracted_name, "span": (0, 0)})
        print(f"[INFO] ✓ 정규식 이름 탐지: {extracted_name}")
    
    # 화이트리스트 이름 탐지 (NER이 놓치는 흔한 이름)
    for whitelist_name in NAME_WHITELIST:
        if whitelist_name in Text and whitelist_name not in [d['value'] for d in Detected]:
            start_idx = Text.find(whitelist_name)
            Detected.append({"type": "PS", "value": whitelist_name, "span": (start_idx, start_idx + len(whitelist_name))})
            print(f"[INFO] ✓ 화이트리스트 이름 탐지 (PS): {whitelist_name}")
    
    for Entity in Results:
        Label = Entity['entity_group']
        Word = Entity['word']
        Start = Entity.get('start')
        End = Entity.get('end')
        
        # ## 접두사 제거 (BERT 토큰화 부산물)
        if Word.startswith('##'):
            Word = Word[2:]
        
        print(f"[DEBUG] 엔티티: Label={Label}, Word={Word}")
        
        if Start is None or End is None:
            continue
        
        if Label.upper() in ["PER", "PS", "ORG", "OG", "LC", "DT"]:
            # LC (주소) 병합 처리
            if Label.upper() == "LC":
                location_parts.append(Word.strip())
                continue
            
            # ORG/OG (조직명) 처리
            if Label.upper() in ["ORG", "OG"]:
                clean_org = Word.replace(" ", "").strip()
                if len(clean_org) >= 2:
                    # 조직명 분리: "삼성전자 개발팀" → "삼성전자", "개발팀"
                    org_split_keywords = ['팀', '부', '부서', '본부', '지점', '센터', '연구소']
                    split_orgs = []
                    
                    # 공백으로 분리 먼저 시도
                    if ' ' in Word:
                        parts = Word.split(' ', 1)
                        if len(parts) == 2:
                            split_orgs = [parts[0].strip(), parts[1].strip()]
                            print(f"[INFO] ⚠ 조직명 분리 (공백): '{Word}' → {split_orgs}")
                    
                    # 공백 없으면 키워드로 분리
                    if not split_orgs:
                        for keyword in org_split_keywords:
                            if keyword in Word and not Word.endswith(keyword):
                                idx = Word.find(keyword)
                                if idx > 0:
                                    part1 = Word[:idx].strip()
                                    part2 = Word[idx:].strip()
                                    if part1 and part2:
                                        split_orgs = [part1, part2]
                                        print(f"[INFO] ⚠ 조직명 분리 (키워드): '{Word}' → {split_orgs}")
                                        break
                    
                    if split_orgs:
                        for org in split_orgs:
                            if len(org.replace(" ", "")) >= 2:
                                Detected.append({"type": "ORG", "value": org, "span": (Start, End)})
                                print(f"[INFO] ✓ NER 탐지 (ORG 분리): {org}")
                    else:
                        Detected.append({"type": "ORG", "value": Word, "span": (Start, End)})
                        print(f"[INFO] ✓ NER 탐지 (ORG): {Word}")
                else:
                    print(f"[INFO] ✗ NER 필터링 (ORG): {Word} (너무 짧음)")
                continue
            
            # DT (날짜/숫자) - 학번/사번으로 처리
            if Label.upper() == "DT":
                # 숫자만 있고 8자리 이상이면 학번/사번으로 간주
                if Word.isdigit() and len(Word) >= 8:
                    Detected.append({"type": "student_id", "value": Word, "span": (Start, End)})
                    print(f"[INFO] ✓ NER 탐지 (student_id): {Word}")
                continue
            
            # PS (사람 이름) 필터링
            if Label.upper() in ["PS", "PER"]:
                # 특수문자 제거 (앞뒤만)
                Word = re.sub(r'^[^가-힣a-zA-Z]+', '', Word)
                Word = re.sub(r'[^가-힣a-zA-Z]+$', '', Word)
                Word = Word.strip()
                
                if not Word:
                    continue
                
                clean_word = Word.replace(" ", "").strip()
                
                # 한글 포함 여부 확인
                has_korean = any('\uac00' <= c <= '\ud7a3' for c in Word)
                if not has_korean:
                    print(f"[INFO] ✗ NER 필터링 ({Label}): {Word} (한글 없음)")
                    continue
                
                # 조직명 키워드 체크 (PS로 잘못 탐지된 경우)
                org_keywords = ['회사', '전자', '그룹', '기업', '주식회사', '(주)', '㈜', 
                               '학교', '대학교', '대학', '고등학교', '중학교', '초등학교',
                               '병원', '의원', '센터', '연구소', '재단', '협회', '은행',
                               '부서', '팀', '본부', '지점', '영업소']
                if any(keyword in Word for keyword in org_keywords):
                    print(f"[INFO] ✗ NER 필터링 ({Label}): {Word} (조직명 키워드) → ORG로 재분류")
                    Detected.append({"type": "ORG", "value": Word, "span": (Start, End)})
                    continue
                
                # 신뢰도 기반 필터링
                score = Entity.get('score', 0)
                
                # 2글자: 신뢰도 90% 이상
                if len(clean_word) == 2 and score < 0.90:
                    print(f"[INFO] ✗ NER 필터링 ({Label}): {Word} (2글자, 신뢰도 {score:.2f} < 0.90)")
                    continue
                
                # 4글자 이상: 신뢰도 80% 이상
                if len(clean_word) >= 4 and score < 0.80:
                    print(f"[INFO] ✗ NER 필터링 ({Label}): {Word} (4글자+, 신뢰도 {score:.2f} < 0.80)")
                    continue
                
                # 숫자만 있는 경우 제외
                if clean_word.isdigit():
                    print(f"[INFO] ✗ NER 필터링 ({Label}): {Word} (숫자만 포함)")
                    continue
            
            Detected.append({"type": Label, "value": Word, "span": (Start, End)})
            print(f"[INFO] ✓ NER 탐지 ({Label}): {Word}")
    
    # 주소 병합
    if location_parts:
        merged_location = ' '.join(location_parts).strip()
        if merged_location:
            Detected.append({"type": "LC", "value": merged_location, "span": (0, 0)})
            print(f"[INFO] ✓ 주소 병합: {merged_location}")
    
    return Detected

# ==========================
# 준식별자 패턴 탐지
# ==========================
def detect_quasi_identifiers(text: str) -> list:
    """학번/사번 등 준식별자 패턴 탐지"""
    detected = []
    
    # 학번/사번 패턴 (8~10자리 숫자, 19xx 또는 20xx로 시작)
    # 여권번호와 겹치지 않도록 숫자만 확인
    id_pattern = re.compile(r'\b(19|20)\d{6,8}\b')
    for match in id_pattern.finditer(text):
        value = match.group()
        # 여권번호 패턴 제외 (알파벳 포함 여부 확인)
        start, end = match.span()
        # 앞뒤에 알파벳이 있으면 여권번호일 가능성
        if start > 0 and text[start-1].isalpha():
            continue
        if end < len(text) and text[end].isalpha():
            continue
        
        detected.append({
            "type": "student_id",
            "value": value,
            "span": match.span()
        })
    
    return detected

# ==========================
# 카테고리 분류 및 위험도 분석
# ==========================
def categorize_detection(item):
    """탐지 항목을 카테고리로 분류"""
    item_type = item.get('type', '')
    
    # 식별자 (개인을 직접 특정 가능)
    if item_type in ['PS', 'PER', 'ssn', 'phone', 'email', 'card', 'account', 'passport', 'driver_license']:
        return 'identifier'
    
    # 민감정보 (얼굴사진)
    if item_type in ['image_face']:
        return 'sensitive'
    
    # 준식별자 (조합시 특정 가능)
    if item_type in ['ORG', 'OG', 'student_id', 'birth', 'LC']:
        return 'quasi'
    
    return 'other'

def analyze_combination_risk(detected_items, text):
    """조합 위험도 자동 분석"""
    print(f"[DEBUG] 조합 분석 시작: 총 {len(detected_items)}개 항목")
    
    if len(detected_items) < 2:
        print(f"[DEBUG] 조합 분석 스킵: 항목이 2개 미만 ({len(detected_items)}개)")
        return None
    
    # 카테고리별 분류
    categorized = {}
    for item in detected_items:
        category = categorize_detection(item)
        if category not in categorized:
            categorized[category] = []
        categorized[category].append(item)
        print(f"[DEBUG] 항목 분류: type={item.get('type')}, value={item.get('value')[:20] if item.get('value') else 'N/A'} → category={category}")
    
    identifier_count = len(categorized.get('identifier', []))
    quasi_count = len(categorized.get('quasi', []))
    sensitive_count = len(categorized.get('sensitive', []))
    
    print(f"[DEBUG] 조합 분석 결과: 식별자={identifier_count}, 준식별자={quasi_count}, 민감정보={sensitive_count}")
    print(f"[DEBUG] 준식별자 항목: {[item.get('type') for item in categorized.get('quasi', [])]}") 
    
    # 위험도 계산
    risk_level = None
    risk_message = None
    risk_items = []
    
    if sensitive_count > 0 and identifier_count > 0:
        risk_level = 'critical'
        risk_message = '민감정보 + 식별자 조합 → 개인 완전 특정 가능!'
        risk_items = categorized.get('identifier', []) + categorized.get('sensitive', [])
    elif identifier_count >= 3:
        risk_level = 'high'
        risk_message = f'식별자 {identifier_count}개 → 개인 특정 가능'
        risk_items = categorized.get('identifier', [])
    elif identifier_count >= 2:
        risk_level = 'medium'
        risk_message = f'식별자 {identifier_count}개 → 개인 특정 가능성 있음'
        risk_items = categorized.get('identifier', [])
    elif identifier_count >= 1 and quasi_count >= 2:
        risk_level = 'high'
        risk_message = f'식별자 + 준식별자 {quasi_count}개 조합 → 개인 특정 가능'
        risk_items = categorized.get('identifier', []) + categorized.get('quasi', [])
    elif identifier_count >= 1 and quasi_count >= 1:
        risk_level = 'medium'
        risk_message = '식별자 + 준식별자 조합 → 개인 특정 가능성 있음'
        risk_items = categorized.get('identifier', []) + categorized.get('quasi', [])
    elif quasi_count >= 2:
        quasi_items = categorized.get('quasi', [])
        non_org_quasi = [item for item in quasi_items if item.get('type') not in ['ORG', 'OG']]
        if len(non_org_quasi) >= 1:
            risk_level = 'medium'
            risk_message = f'준식별자 {quasi_count}개 조합 → 개인 특정 가능성 있음'
            risk_items = quasi_items
        else:
            print(f"[DEBUG] 조합 위험 스킵: 조직명만 {len(quasi_items)}개 (개인 특정 불가)")
    
    if risk_level:
        print(f"[WARN] ⚠️ 조합 위험 감지: {risk_level} - {risk_message}")
        return {
            'level': risk_level,
            'message': risk_message,
            'items': risk_items,
            'counts': {
                'identifier': identifier_count,
                'quasi': quasi_count,
                'sensitive': sensitive_count
            }
        }
    
    return None

# ==========================
# 얼굴 탐지 (MTCNN)
# ==========================
def detect_faces_in_image_bytes(image_bytes, confidence_threshold=0.98):
    """이미지에서 얼굴을 탐지합니다 (오탐지 방지 강화)"""
    try:
        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        
        # 이미지 크기 체크 (너무 작으면 스킵)
        if img.width < 50 or img.height < 50:
            return []
        
        # 이미지 리사이즈 (큰 이미지는 축소하여 속도 향상)
        max_size = 800
        if img.width > max_size or img.height > max_size:
            ratio = min(max_size / img.width, max_size / img.height)
            new_size = (int(img.width * ratio), int(img.height * ratio))
            img = img.resize(new_size, Image.LANCZOS)
        
        img_np = np.array(img)
        results = detector.detect_faces(img_np)
        detections = []
        
        for res in results:
            conf = float(res['confidence'])
            x, y, w, h = res['box']
            
            # 필터링 조건
            # 1. confidence >= 0.98 (더 엄격하게)
            if conf < confidence_threshold:
                continue
            
            # 2. 얼굴 크기 검증 (너무 작거나 큰 것 제외)
            if w < 30 or h < 30 or w > img.width * 0.9 or h > img.height * 0.9:
                continue
            
            # 3. 가로세로 비율 검증 (얼굴은 대략 0.7~1.3 비율)
            aspect_ratio = w / h if h > 0 else 0
            if aspect_ratio < 0.6 or aspect_ratio > 1.5:
                continue
            
            # 4. keypoints 검증 (눈, 코, 입이 제대로 탐지되었는지)
            keypoints = res.get('keypoints', {})
            if keypoints:
                left_eye = keypoints.get('left_eye')
                right_eye = keypoints.get('right_eye')
                nose = keypoints.get('nose')
                
                # 눈과 코가 모두 탐지되어야 함
                if not (left_eye and right_eye and nose):
                    continue
                
                # 두 눈 사이 거리 검증
                eye_distance = abs(left_eye[0] - right_eye[0])
                if eye_distance < w * 0.2 or eye_distance > w * 0.8:
                    continue
            
            detections.append({
                "bbox": [int(x), int(y), int(w), int(h)], 
                "confidence": conf
            })
        
        if len(detections) > 0:
            print(f"[INFO] ✓ MTCNN 얼굴 탐지 성공: {len(detections)}개 발견 (필터링 후)")
        return detections
    except Exception as e:
        print(f"[ERROR] MTCNN 얼굴 탐지 실패: {e}")
        return []

# ==========================
# DOCX / PDF 이미지 탐지
# ==========================
def detect_images_in_docx(file_bytes):
    results = []
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as docx_zip:
            image_files = [f for f in docx_zip.namelist() if f.startswith("word/media/")]
            print(f"[INFO] DOCX 내부 이미지 {len(image_files)}개 발견, 얼굴 탐지 시작")
            for image_name in image_files:
                with docx_zip.open(image_name) as image_file:
                    img_bytes = image_file.read()
                    faces = detect_faces_in_image_bytes(img_bytes)
                    if len(faces) > 0:
                        print(f"[INFO] ✓ {image_name} → 얼굴 {len(faces)}개 탐지")
                        results.append({
                            "image_name": image_name,
                            "faces_found": len(faces),
                            "faces": faces
                        })
    except Exception as e:
        print(f"[WARN] DOCX 이미지 검사 실패: {e}")
    return results

def _process_pdf_image(args):
    """PDF 이미지 하나를 처리 (병렬 처리용)"""
    pno, img_index, img_bytes = args
    faces = detect_faces_in_image_bytes(img_bytes)
    if len(faces) > 0:
        print(f"[INFO] ✓ PDF p{pno+1}, img{img_index} → 얼굴 {len(faces)}개 탐지")
        return {
            "page": pno + 1,
            "image_index": img_index,
            "faces_found": len(faces),
            "faces": faces
        }
    return None

def detect_images_in_pdf(pdf_bytes):
    results = []
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            # 중복 이미지 제거 (xref 기반)
            seen_xrefs = set()
            image_tasks = []
            
            for pno, page in enumerate(doc):
                images = page.get_images(full=True)
                if len(images) > 0:
                    print(f"[INFO] PDF 페이지 {pno+1}에서 이미지 {len(images)}개 발견")
                
                for img_index, img in enumerate(images):
                    xref = img[0]
                    
                    # 중복 이미지 스킵 (같은 이미지가 여러 페이지에 반복되는 경우)
                    if xref in seen_xrefs:
                        continue
                    seen_xrefs.add(xref)
                    
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image.get("image")
                    
                    # 너무 작은 이미지 스킵 (로고, 아이콘 등)
                    if img_bytes and len(img_bytes) > 5000:  # 5KB 이상만
                        image_tasks.append((pno, img_index, img_bytes))
            
            # 병렬 처리 (CPU 코어 수에 맞춰 조정)
            if image_tasks:
                print(f"[INFO] 총 {len(image_tasks)}개 이미지를 병렬 처리 시작 (중복 제거 완료)")
                max_workers = min(8, os.cpu_count() or 4)  # CPU 코어 수에 맞춰 동적 조정
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    futures = [executor.submit(_process_pdf_image, task) for task in image_tasks]
                    for future in as_completed(futures):
                        result = future.result()
                        if result:
                            results.append(result)
    except Exception as e:
        print(f"[WARN] PDF 이미지 검사 실패: {e}")
    return results

def detect_images_in_hwp(hwp_bytes):
    """HWP 파일 내부 이미지에서 얼굴 탐지"""
    results = []
    if olefile is None:
        return results
    try:
        ole = olefile.OleFileIO(io.BytesIO(hwp_bytes))
        for entry in ole.listdir():
            if entry[0] == "BinData":
                try:
                    stream = ole.openstream(entry)
                    img_bytes = stream.read()
                    faces = detect_faces_in_image_bytes(img_bytes)
                    if len(faces) > 0:
                        img_name = "/".join(entry)
                        print(f"[INFO] ✓ {img_name} → 얼굴 {len(faces)}개 탐지")
                        results.append({
                            "image_name": img_name,
                            "faces_found": len(faces),
                            "faces": faces
                        })
                except Exception:
                    continue
        ole.close()
    except Exception as e:
        print(f"[WARN] HWP 이미지 검사 실패: {e}")
    return results

def detect_images_in_pptx(pptx_bytes):
    """PPTX 파일 내부 이미지에서 얼굴 탐지"""
    results = []
    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as pptx_zip:
            image_files = [f for f in pptx_zip.namelist() if f.startswith("ppt/media/")]
            print(f"[INFO] PPTX 내부 이미지 {len(image_files)}개 발견, 얼굴 탐지 시작")
            for image_name in image_files:
                with pptx_zip.open(image_name) as image_file:
                    img_bytes = image_file.read()
                    faces = detect_faces_in_image_bytes(img_bytes)
                    if len(faces) > 0:
                        print(f"[INFO] ✓ {image_name} → 얼굴 {len(faces)}개 탐지")
                        results.append({
                            "image_name": image_name,
                            "faces_found": len(faces),
                            "faces": faces
                        })
    except Exception as e:
        print(f"[WARN] PPTX 이미지 검사 실패: {e}")
    return results

def scan_file_for_face_images(file_bytes, file_ext):
    """파일 내 이미지에서 얼굴을 탐지합니다"""
    file_ext = file_ext.lower()
    print(f"[INFO] === 얼굴 이미지 스캔 시작 ({file_ext}) ===")
    
    if file_ext == "docx":
        return detect_images_in_docx(file_bytes)
    elif file_ext == "pdf":
        return detect_images_in_pdf(file_bytes)
    elif file_ext == "hwp":
        return detect_images_in_hwp(file_bytes)
    elif file_ext == "pptx":
        return detect_images_in_pptx(file_bytes)
    elif file_ext in ["png", "jpg", "jpeg", "bmp", "webp", "gif", "tiff"]:
        print(f"[INFO] 단일 이미지 파일 얼굴 탐지 시작")
        faces = detect_faces_in_image_bytes(file_bytes)
        if len(faces) > 0:
            return [{
                "image_name": "uploaded_image", 
                "faces_found": len(faces), 
                "faces": faces
            }]
        return []
    else:
        return []



# ==========================
# 최종 핸들러 (마스킹 제거)
# ==========================
def handle_input_raw(Input_Data, Original_Format=None, Original_Filename=None):
    """
    파일을 처리하는 메인 함수
    
    Returns:
        - detected: 탐지된 개인정보 리스트
        - masked_text: 빈 문자열 (마스킹 제거됨)
        - backend_status: 백엔드 전송 성공 여부
        - image_detections: 이미지 내 얼굴 탐지 결과
    """
    if not isinstance(Input_Data, bytes):
        raise ValueError("지원하지 않는 입력 형식입니다.")
    
    print(f"\n[INFO] ========== 파일 처리 시작 (확장자: {Original_Format}) ==========")
    
    # 1단계: 파일 파싱 및 텍스트 추출
    try:
        Parsed_Text, is_pure_image = parse_file(Input_Data, Original_Format)
        Parsed_Text = Parsed_Text.strip()
    except ValueError as ve:
        raise ve
    
    # 2단계: 이미지 내 얼굴 탐지
    image_detections = scan_file_for_face_images(Input_Data, Original_Format or "")
    
    # 3단계: 텍스트에서 개인정보 탐지 (파일명 포함)
    Detected = []
    
    # 파일명을 텍스트에 포함
    combined_text = Parsed_Text
    if Original_Filename:
        filename_without_ext = Original_Filename.rsplit('.', 1)[0]
        combined_text = filename_without_ext + " " + Parsed_Text
    
    if combined_text.strip():
        print(f"[INFO] 텍스트 분석 중... (길이: {len(combined_text)} 글자)")
        print(f"[DEBUG] 텍스트 샘플: {combined_text[:500]}")
        
        # NER 먼저 실행
        ner_results = detect_by_ner(combined_text)
        
        # 정규식 실행
        regex_results = detect_by_regex(combined_text)
        
        # 준식별자 패턴 탐지
        quasi_results = detect_quasi_identifiers(combined_text)
        
        # 중복 제거: value 기준으로 중복 제거
        all_detected = []
        seen = set()
        for item in regex_results + ner_results + quasi_results:
            val = item['value'].strip().lower()
            if val not in seen:
                seen.add(val)
                all_detected.append(item)
        
        print(f"[INFO] ✓ 텍스트 개인정보 {len(all_detected)}개 탐지 (중복 제거 완료)")
        
        # 조합 위험도 분석 (모든 항목 포함)
        combination_risk = analyze_combination_risk(all_detected, combined_text)
        
        # 조합 위험도 분석 결과에 따라 처리
        if combination_risk:
            # 조합 위험도 있으면 모든 항목 포함
            Detected = all_detected
            print(f"[WARN] ⚠️ 조합 위험 감지: {combination_risk['level']} - {combination_risk['message']}")
            Detected.append({
                "type": "combination_risk",
                "value": combination_risk['message'],
                "risk_level": combination_risk['level'],
                "risk_items": combination_risk['items'],
                "counts": combination_risk['counts']
            })
        else:
            # 조합 위험도 없으면: 식별자/민감정보만 포함 (준식별자 제외)
            Detected = [item for item in all_detected if categorize_detection(item) in ['identifier', 'sensitive']]
    else:
        print("[INFO] 추출된 텍스트 없음")
    
    # 4단계: 얼굴 탐지 결과를 Detected에 추가
    total_faces = 0
    for img in image_detections:
        face_count = img.get("faces_found", 0)
        total_faces += face_count
        if face_count > 0:
            img_name = img.get("image_name", "이미지")
            Detected.append({
                "type": "image_face",
                "value": f"{img_name} 내 얼굴 {face_count}개",
                "detail": img
            })
    
    if total_faces > 0:
        print(f"[INFO] ✓ 이미지 얼굴 총 {total_faces}개 탐지")
    
    # 5단계: 파일명 마스킹
    masked_filename = ""
    if Original_Filename:
        parts = Original_Filename.rsplit('.', 1)
        name_part = parts[0]
        ext_part = f".{parts[1]}" if len(parts) > 1 else ""
        
        type_map = {
            'PS': '이름', 'PER': '이름',
            'phone': '전화번호',
            'email': '이메일',
            'ssn': '주민번호',
            'birth': '생년월일',
            'card': '카드번호',
            'account': '계좌번호',
            'passport': '여권번호',
            'driver_license': '면허번호'
        }
        
        replacements = []
        for item in Detected:
            value = item.get('value', '').strip()
            item_type = item.get('type', '')
            if value and value in name_part and item_type in type_map:
                replacements.append((value, type_map[item_type]))
        
        replacements.sort(key=lambda x: len(x[0]), reverse=True)
        
        for original, type_name in replacements:
            name_part = name_part.replace(original, f"({type_name})")
        
        masked_filename = name_part + ext_part
        if masked_filename != Original_Filename:
            print(f"[INFO] 파일명 변경: {Original_Filename} → {masked_filename}")
    
    # 6단계: 처리 완료
    if not Detected:
        print("[INFO] ========== 탐지된 민감정보 없음 ==========\n")
    else:
        print(f"[INFO] ========== 처리 완료 - 민감정보: {len(Detected)}개 ==========\n")
    
    return Detected, masked_filename, False, image_detections